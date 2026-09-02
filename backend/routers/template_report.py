"""Template Report — reusable ChartBuilder report layouts and PPTX export.

한 장짜리 차트 모음이 아니라 "매번 같은 형식으로 다시 뽑는 보고서"가 목표다.
그래서 두 가지를 템플릿이 갖는다.

1. **변수** — 저장된 ChartBuilder 코드 안의 ``{{LOT}}`` 같은 토큰(core/report_variables).
   랏이 바뀌어도 차트를 새로 만들지 않는다.
2. **반복** — 변수 하나(기본 ``LOT``)에 값을 여러 개 주면 페이지 묶음이 랏마다 반복된다.

기본 실행은 ``RECENT_DAYS = 7``처럼 저장 차트가 가진 조건을 그대로 쓴다. 다만 보고서
실행 컨텍스트를 켜면 root lot/wafer 목록과 시간 열·최근 일수를 모든 차트에 한 번만
덧씌울 수 있다. 컨텍스트는 Template 자체를 변경하지 않는 실행 전용 값이다.

A/B 비교는 템플릿이 슬롯을 복제하는 방식이 아니라 **ChartBuilder 코드 쪽에서** 만든다
(조건 열을 COLOR/X 로 잡거나 조건별 차트를 따로 저장). ``stats`` 블록은 그렇게 만들어진
차트 하나를 가리켜 그룹별 통계를 표로 깔고, 그룹이 정확히 둘이면 Δ 와 Δ% 를 덧붙인다.

슬라이드 모양은 HOL Auto Report(`auto report/My_Function.py`)와 같은 규약을 쓴다 —
상단 네이비 바 + 흰 제목, 표지 슬라이드, 하단 푸터. 색은 RGB(31,73,125) 한 값이다.
"""
from __future__ import annotations

import base64
import datetime as dt
import io
import json
import math
import re
import threading
import uuid
import zipfile
from pathlib import Path
from urllib.parse import quote

from fastapi import APIRouter, Depends, HTTPException
from fastapi.responses import Response
from pydantic import BaseModel

from core.auth import current_user, require_page_manager
from core.chart_builder_definition import ChartBuilderDefinitionError, linked_chart_color_pairs, parse_chart_builder_definition
from core.paths import PATHS
from core.report_variables import (
    ReportVariableError,
    extract_variables,
    normalize_name,
    split_list,
    substitute,
    validate_bindings,
)
from core.utils import load_json, save_json, safe_filename


router = APIRouter(prefix="/api/template-report", tags=["template-report"])
STORE_FILE = PATHS.data_root / "template_reports.json"
SETTINGS_FILE = PATHS.data_root / "template_report_settings.json"
BACKGROUND_FILE = PATHS.data_root / "template_report_background.png"
_STORE_LOCK = threading.Lock()
_SETTINGS_LOCK = threading.Lock()
MAX_PAGES = 30
MAX_CHARTS_PER_PAGE = 20
MAX_REPEAT_VALUES = 20
MAX_RENDERED_PAGES = 200
MAX_IMAGES = MAX_RENDERED_PAGES * MAX_CHARTS_PER_PAGE
LEGACY_SLOT_LAYOUTS = {
    1: {"x": 3.4, "y": 13.6, "width": 45.2, "height": 37.1},
    2: {"x": 51.5, "y": 13.6, "width": 45.2, "height": 37.1},
    3: {"x": 3.4, "y": 54.4, "width": 45.2, "height": 37.1},
    4: {"x": 51.5, "y": 54.4, "width": 45.2, "height": 37.1},
}
DEFAULT_SLOT_LAYOUT = {"x": 5.0, "y": 16.0}
SLIDE_DESIGN_WIDTH = 1920
SLIDE_DESIGN_HEIGHT = 1080
DEFAULT_CHART_WIDTH = 1200
DEFAULT_CHART_HEIGHT = 650
# 차트가 아닌 블록(표·글)은 ChartBuilder 크기가 없어 슬라이드 비율로 직접 잡는다.
DEFAULT_BLOCK_WIDTH_PCT = 46.0
DEFAULT_BLOCK_HEIGHT_PCT = 26.0
SLOT_KINDS = ("chart", "split", "text", "stats", "legend")
# 예전 템플릿의 A/B 전용 통계표 — 지금은 그룹 통계표 하나로 합쳤다.
LEGACY_SLOT_KINDS = {"ab_stats": "stats"}

# ─ 슬라이드 디자인 상수 — HOL Auto Report 와 같은 값 ──────────────────────────
SLIDE_WIDTH_IN = 13.333333
SLIDE_HEIGHT_IN = 7.5
TITLE_BAR_HEIGHT_IN = 0.42
COVER_BAR_HEIGHT_IN = 0.45
# Flow Web의 IBM Carbon 레이어와 같은 중립 면 + 오렌지 포인트.
REPORT_ACCENT = (226, 88, 34)
REPORT_TEXT = (23, 23, 23)
REPORT_MUTED = (115, 115, 115)
REPORT_PAGE = (250, 250, 250)
REPORT_PANEL = (255, 255, 255)
REPORT_SUBTLE = (245, 245, 245)
REPORT_BORDER = (229, 229, 229)
REPORT_BORDER_STRONG = (163, 163, 163)
REPORT_FONT = "Malgun Gothic"
MAX_TABLE_ROWS = 26
MAX_TABLE_COLUMNS = 24
MAX_BACKGROUND_BYTES = 12 * 1024 * 1024
MAX_BACKGROUND_PIXELS = 40_000_000
DEFAULT_SETTINGS = {
    # 배경 파일이 없으면 기존 Carbon 중립 배경을 그대로 쓴다. config seed가 없어도
    # 첫 실행과 setup.py 배포에서 동작해야 하므로 기본값은 코드가 소유한다.
    "background_updated_by": "",
    "background_updated_at": "",
}


class TemplateSlotReq(BaseModel):
    position: int
    chart_id: str = ""
    kind: str = "chart"
    x: float | None = None
    y: float | None = None
    width: float | None = None
    height: float | None = None
    chart_width: int | None = None
    chart_height: int | None = None
    chart_name: str = ""
    # Template 안에서 저장 차트의 생성식을 복사해 직접 수정할 수 있다.
    # chart_id는 원본/계보를 가리키고 definition_code는 이 슬롯이 실제 실행할 코드다.
    definition_code: str = ""
    title: str = ""
    # kind=text
    text: str = ""
    # kind=split
    product: str = ""
    lot: str = ""
    columns: str = ""
    display_mode: str = "matrix"
    # kind=stats / legend — 같은 페이지의 chart position을 가리킨다.
    source_position: int = 0
    stats: str = ""


class TemplatePageReq(BaseModel):
    id: str = ""
    title: str = ""
    subtitle: str = ""
    slots: list[TemplateSlotReq] = []


class TemplateVariableReq(BaseModel):
    name: str
    label: str = ""
    default: str = ""


class TemplateOptionsReq(BaseModel):
    cover: bool = True
    footer: bool = True
    subtitle: str = ""
    repeat_variable: str = "LOT"


class TemplateSaveReq(BaseModel):
    id: str = ""
    name: str
    pages: list[TemplatePageReq]
    variables: list[TemplateVariableReq] = []
    options: TemplateOptionsReq | None = None


class TemplateCodeReq(BaseModel):
    code: str


class TemplateAssistantReq(BaseModel):
    instruction: str
    template_code: str


class TemplateRunContextReq(BaseModel):
    root_lot_ids: list[str] = []
    wafer_ids: list[str] = []
    override_recent_days: bool = False
    recent_days: int = 0
    date_column: str = "tkout_time"
    color_rules: list[str] = []
    color_else: str = "gray"


class TemplateRunReq(BaseModel):
    template_id: str
    bindings: dict[str, str] = {}
    repeat_values: list[str] = []
    context: TemplateRunContextReq | None = None


class SplitBlockReq(BaseModel):
    product: str
    lot_id: str
    columns: str = ""
    display_mode: str = "matrix"
    max_rows: int = MAX_TABLE_ROWS


class ExportImageReq(BaseModel):
    page_index: int = 0
    position: int = 0
    key: str = ""
    chart_id: str = ""
    data_url: str


class ExportTableReq(BaseModel):
    key: str
    title: str = ""
    columns: list[str] = []
    rows: list[list[str]] = []
    note: str = ""


class ExportReq(BaseModel):
    template_id: str
    bindings: dict[str, str] = {}
    repeat_values: list[str] = []
    context: TemplateRunContextReq | None = None
    images: list[ExportImageReq] = []
    tables: list[ExportTableReq] = []


class BackgroundSaveReq(BaseModel):
    data_url: str


# ── 저장소 ────────────────────────────────────────────────────────────────────
def _load_templates() -> list[dict]:
    raw = load_json(STORE_FILE, {}) or {}
    rows = raw.get("templates") if isinstance(raw, dict) else []
    return _normalize_template_names([row for row in (rows or []) if isinstance(row, dict)])


def _unique_template_name(base: str, used: set[str]) -> str:
    clean = _clean_text(base, 120) or "Template Report"
    candidate = clean
    number = 2
    while candidate.casefold() in used:
        suffix = f" ({number})"
        candidate = f"{clean[:max(1, 120 - len(suffix))]}{suffix}"
        number += 1
    return candidate


def _normalize_template_names(rows: list[dict]) -> list[dict]:
    used = {
        str(row.get("id") or "").strip().casefold()
        for row in rows
        if row.get("id")
    }
    normalized = []
    for raw in rows:
        row = dict(raw)
        name = _unique_template_name(str(row.get("name") or "Template Report"), used)
        row["name"] = name
        used.add(name.casefold())
        normalized.append(row)
    return normalized


def _save_templates(rows: list[dict]) -> None:
    save_json(STORE_FILE, {"version": 2, "templates": rows})


def _load_settings() -> dict:
    raw = load_json(SETTINGS_FILE, {}) or {}
    if not isinstance(raw, dict):
        raw = {}
    return {**DEFAULT_SETTINGS, **raw}


def _background_public_settings(*, include_data: bool = False) -> dict:
    settings = _load_settings()
    configured = BACKGROUND_FILE.is_file()
    version = BACKGROUND_FILE.stat().st_mtime_ns if configured else 0
    result = {
        "configured": configured,
        "url": f"/api/template-report/settings/background/image?v={version}" if configured else "",
        "updated_by": _clean_text(settings.get("background_updated_by"), 120),
        "updated_at": _clean_text(settings.get("background_updated_at"), 80),
    }
    if include_data:
        payload = _background_bytes() if configured else b""
        result["data_url"] = f"data:image/png;base64,{base64.b64encode(payload).decode('ascii')}" if payload else ""
    return result


def _decode_background(data_url: str) -> bytes:
    """브라우저 clipboard 이미지를 PPTX가 안정적으로 읽는 PNG로 정규화한다."""
    match = re.fullmatch(
        r"data:image/(png|jpeg|jpg|webp);base64,([A-Za-z0-9+/=\s]+)",
        str(data_url or ""),
        re.IGNORECASE,
    )
    if not match:
        raise HTTPException(400, "배경은 PNG, JPEG 또는 WebP 이미지여야 합니다.")
    try:
        payload = base64.b64decode(match.group(2), validate=False)
    except Exception as exc:
        raise HTTPException(400, "배경 이미지 base64를 해석하지 못했습니다.") from exc
    if not payload or len(payload) > MAX_BACKGROUND_BYTES:
        raise HTTPException(400, "배경 이미지는 12MB 이하여야 합니다.")

    try:
        from PIL import Image, UnidentifiedImageError

        with Image.open(io.BytesIO(payload)) as source:
            source.load()
            if source.width < 1 or source.height < 1 or source.width * source.height > MAX_BACKGROUND_PIXELS:
                raise HTTPException(400, "배경 이미지 해상도가 너무 큽니다. 최대 4천만 픽셀까지 사용할 수 있습니다.")
            image = source.convert("RGBA" if "A" in source.getbands() else "RGB")
            out = io.BytesIO()
            image.save(out, format="PNG", optimize=True)
            normalized = out.getvalue()
    except HTTPException:
        raise
    except (UnidentifiedImageError, OSError, ValueError) as exc:
        raise HTTPException(400, "유효한 배경 이미지가 아닙니다.") from exc
    if len(normalized) > MAX_BACKGROUND_BYTES:
        raise HTTPException(400, "PNG로 변환된 배경 이미지가 12MB를 넘습니다.")
    return normalized


def _background_bytes() -> bytes:
    try:
        payload = BACKGROUND_FILE.read_bytes()
    except OSError:
        return b""
    return payload if payload else b""


def _chart_history() -> dict[str, dict]:
    from routers import filebrowser

    # ChartBuilder와 같은 카탈로그 계약: Auto Report/관리자 고정 전체 + 최근 500.
    entries = filebrowser._chart_builder_visible_history_entries(recent_limit=500)
    return {str(entry.get("history_id") or ""): entry for entry in entries if entry.get("history_id")}


def _chart_history_by_name(history: dict[str, dict]) -> dict[str, dict]:
    return {
        str(entry.get("name") or "").strip().casefold(): entry
        for entry in history.values()
        if str(entry.get("name") or "").strip()
    }


def _template_or_404(template_id: str) -> dict:
    template_ref = str(template_id or "").strip()
    rows = _load_templates()
    row = next((item for item in rows if str(item.get("id")) == template_ref), None)
    if not row:
        folded = template_ref.casefold()
        row = next((item for item in rows if str(item.get("name") or "").casefold() == folded), None)
    if not row:
        raise HTTPException(404, "Template Report를 찾지 못했습니다.")
    return row


def _template_options(row: dict) -> dict:
    raw = row.get("options") if isinstance(row.get("options"), dict) else {}
    repeat_variable = normalize_name(raw.get("repeat_variable", "LOT"))
    return {
        "cover": bool(raw.get("cover", True)),
        "footer": bool(raw.get("footer", True)),
        "subtitle": _clean_text(raw.get("subtitle", ""), 180),
        "repeat_variable": repeat_variable,
    }


def _template_variables(row: dict) -> list[dict]:
    declared = row.get("variables") if isinstance(row.get("variables"), list) else []
    out: list[dict] = []
    seen: set[str] = set()
    for item in declared:
        if not isinstance(item, dict):
            continue
        name = normalize_name(item.get("name"))
        if not name or name in seen:
            continue
        seen.add(name)
        out.append({
            "name": name,
            "label": _clean_text(item.get("label"), 80) or name,
            "default": _clean_text(item.get("default"), 200),
        })
    for name in _used_variable_names(row):
        if name not in seen:
            seen.add(name)
            out.append({"name": name, "label": name, "default": ""})
    return out


def _used_variable_names(row: dict) -> list[str]:
    """템플릿 본문(차트 코드·제목·표 대상)에 실제로 등장하는 변수."""
    texts: list[str] = [str(_template_options(row).get("subtitle") or "")]
    for page in row.get("pages") or []:
        if not isinstance(page, dict):
            continue
        texts.extend([str(page.get("title") or ""), str(page.get("subtitle") or "")])
        for slot in page.get("slots") or []:
            if not isinstance(slot, dict):
                continue
            texts.extend([
                str(slot.get("definition_code") or ""),
                str(slot.get("text") or ""),
                str(slot.get("product") or ""),
                str(slot.get("lot") or ""),
                str(slot.get("columns") or ""),
                str(slot.get("title") or ""),
            ])
    return extract_variables(*texts)


def _public_template(row: dict) -> dict:
    pages = []
    for raw_page in row.get("pages") if isinstance(row.get("pages"), list) else []:
        page = dict(raw_page)
        page["title"] = str(raw_page.get("title") or "")
        page["subtitle"] = str(raw_page.get("subtitle") or "")
        page["slots"] = [
            {**slot, "kind": _slot_kind(slot), **_slot_layout(slot), **_code_time_window(slot.get("definition_code"))}
            for slot in (raw_page.get("slots") or [])
            if isinstance(slot, dict)
        ]
        pages.append(page)
    return {
        "id": str(row.get("id") or ""),
        "name": str(row.get("name") or ""),
        "pages": pages,
        "variables": _template_variables(row),
        "options": _template_options(row),
        "created_by": str(row.get("created_by") or ""),
        "created_at": str(row.get("created_at") or ""),
        "updated_by": str(row.get("updated_by") or ""),
        "updated_at": str(row.get("updated_at") or ""),
        "slide_size": "LAYOUT_WIDE",
    }


RECENT_DAYS_HINT_RE = re.compile(r"(?<![A-Za-z0-9_])recent[_\s-]*days?\s*[:=]\s*(\d+)", re.IGNORECASE)
DATE_COLUMN_HINT_RE = re.compile(r"(?<![A-Za-z0-9_])(?:date|time)[_\s-]*col(?:umn)?\s*[:=]\s*([A-Za-z_]\w*)", re.IGNORECASE)


def _code_time_window(code) -> dict:
    """차트 코드가 스스로 지닌 시간 창 — 화면에 "최근 7일"을 그대로 보여주기 위한 요약.

    실행에 쓰이는 값이 아니라(그건 파서가 낸다) 라벨이라, 코드가 깨져 있어도
    목록이 통째로 죽지 않도록 가벼운 정규식으로만 읽는다.
    """
    text = str(code or "")
    days = [int(match.group(1)) for match in RECENT_DAYS_HINT_RE.finditer(text)]
    if not days:
        return {"recent_days": 0, "date_column": ""}
    columns = [match.group(1) for match in DATE_COLUMN_HINT_RE.finditer(text)]
    return {"recent_days": max(days), "date_column": columns[0] if columns else "tkout_time"}


def _clean_text(value, limit: int) -> str:
    return re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]+", " ", str(value or "")).strip()[:limit]


def _clean_multiline(value, limit: int) -> str:
    text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]+", " ", str(value or "")).strip()
    return text[:limit]


def _slot_kind(slot) -> str:
    raw = slot.get("kind") if isinstance(slot, dict) else getattr(slot, "kind", "")
    kind = str(raw or "chart").strip().casefold()
    kind = LEGACY_SLOT_KINDS.get(kind, kind)
    return kind if kind in SLOT_KINDS else "chart"


def _chart_dimensions(slot: dict | TemplateSlotReq) -> tuple[int, int]:
    if isinstance(slot, dict):
        width = int(slot.get("chart_width") or 0)
        height = int(slot.get("chart_height") or 0)
        definition = str(slot.get("definition_code") or "")
        if (not width or not height) and definition:
            try:
                chart = parse_chart_builder_definition(definition).get("chart") or {}
                width = width or int(chart.get("width") or 0)
                height = height or int(chart.get("height") or 0)
            except (ChartBuilderDefinitionError, TypeError, ValueError):
                pass
    else:
        width = height = 0
    width = max(320, min(2400, width or DEFAULT_CHART_WIDTH))
    height = max(240, min(1600, height or DEFAULT_CHART_HEIGHT))
    scale = min(1.0, SLIDE_DESIGN_WIDTH / width, SLIDE_DESIGN_HEIGHT / height)
    return round(width * scale), round(height * scale)


def _slot_value(slot, key):
    return slot.get(key) if isinstance(slot, dict) else getattr(slot, key, None)


def _slot_layout(slot: dict | TemplateSlotReq) -> dict[str, float | int]:
    """좌상단 좌표 + 블록 크기.

    chart 블록은 기본적으로 ChartBuilder 크기를 쓰되 Template에서 크기를 덮어쓸 수
    있다. 표·글 블록은 슬라이드 비율 입력값을 쓴다.
    """
    kind = _slot_kind(slot)
    position = int(_slot_value(slot, "position") or 1)
    fallback = LEGACY_SLOT_LAYOUTS.get(position, DEFAULT_SLOT_LAYOUT) if kind == "chart" else DEFAULT_SLOT_LAYOUT
    values: dict[str, float | int] = {}
    explicit: dict[str, bool] = {}
    for key in ("x", "y"):
        raw = _slot_value(slot, key)
        explicit[key] = raw is not None
        values[key] = round(float(fallback[key] if raw is None else raw), 3)

    if kind == "chart":
        chart_width, chart_height = _chart_dimensions(slot)
        width_pct = chart_width / SLIDE_DESIGN_WIDTH * 100
        height_pct = chart_height / SLIDE_DESIGN_HEIGHT * 100
    else:
        raw_width = _slot_value(slot, "width")
        raw_height = _slot_value(slot, "height")
        width_pct = float(DEFAULT_BLOCK_WIDTH_PCT if raw_width in (None, 0) else raw_width)
        height_pct = float(DEFAULT_BLOCK_HEIGHT_PCT if raw_height in (None, 0) else raw_height)
        width_pct = min(100.0, max(8.0, width_pct))
        height_pct = min(100.0, max(6.0, height_pct))
        chart_width = round(width_pct / 100 * SLIDE_DESIGN_WIDTH)
        chart_height = round(height_pct / 100 * SLIDE_DESIGN_HEIGHT)

    values.update({
        "chart_width": chart_width,
        "chart_height": chart_height,
        "width": round(width_pct, 3),
        "height": round(height_pct, 3),
    })
    if not all(math.isfinite(float(values[key])) for key in ("x", "y")):
        raise HTTPException(400, "차트 좌표가 올바르지 않습니다.")
    if values["x"] < 0 or values["y"] < 0:
        raise HTTPException(400, "차트 좌표가 올바르지 않습니다.")
    if not explicit["x"]:
        values["x"] = round(min(float(values["x"]), 100 - float(values["width"])), 3)
    if not explicit["y"]:
        values["y"] = round(min(float(values["y"]), 100 - float(values["height"])), 3)
    if values["x"] + values["width"] > 100.001 or values["y"] + values["height"] > 100.001:
        raise HTTPException(400, "차트가 슬라이드 영역을 벗어납니다.")
    return values


# ── 조회 ──────────────────────────────────────────────────────────────────────
@router.get("/templates")
def list_templates(_user=Depends(current_user)):
    rows = sorted(_load_templates(), key=lambda row: str(row.get("updated_at") or ""), reverse=True)
    return {"ok": True, "templates": [_public_template(row) for row in rows]}


@router.get("/charts")
def list_charts(_user=Depends(current_user)):
    rows = list(_chart_history().values())
    charts = []
    for row in rows:
        chart = row.get("chart") if isinstance(row.get("chart"), dict) else {}
        source_text = " · ".join(
            f"{source.get('id')}:{source.get('root')}/{source.get('product')}"
            for source in (row.get("sources") or []) if isinstance(source, dict)
        )
        label = " × ".join(value for value in (str(chart.get("x") or ""), str(chart.get("y") or "")) if value)
        charts.append({
            "id": row.get("history_id"),
            "name": row.get("name"),
            "timestamp": row.get("timestamp"),
            "username": row.get("username"),
            "label": label or source_text or str(row.get("name") or row.get("history_id")),
            "source_text": source_text,
            "chart": chart,
            "definition_code": str(row.get("definition_code") or "")[:100_000],
            "variables": extract_variables(str(row.get("definition_code") or "")),
            "row_count": int(row.get("row_count") or 0),
            "pinned": bool(row.get("pinned")),
            **_code_time_window(row.get("definition_code")),
        })
    return {"ok": True, "charts": charts}


@router.get("/settings")
def get_template_report_settings(_user=Depends(current_user)):
    return {"ok": True, "settings": {"background": _background_public_settings(include_data=True)}}


@router.get("/settings/background/image")
def get_template_report_background(_user=Depends(current_user)):
    payload = _background_bytes()
    if not payload:
        raise HTTPException(404, "설정된 Template Report 배경이 없습니다.")
    return Response(
        content=payload,
        media_type="image/png",
        headers={"Cache-Control": "private, max-age=31536000, immutable"},
    )


@router.post("/settings/background")
def save_template_report_background(
    req: BackgroundSaveReq,
    user=Depends(require_page_manager("templatereport")),
):
    payload = _decode_background(req.data_url)
    now = dt.datetime.now(dt.timezone.utc).isoformat()
    with _SETTINGS_LOCK:
        BACKGROUND_FILE.parent.mkdir(parents=True, exist_ok=True)
        temporary = BACKGROUND_FILE.with_suffix(".tmp")
        temporary.write_bytes(payload)
        temporary.replace(BACKGROUND_FILE)
        settings = _load_settings()
        settings.update({
            "background_updated_by": str(user.get("username") or ""),
            "background_updated_at": now,
        })
        save_json(SETTINGS_FILE, settings)
    return {"ok": True, "settings": {"background": _background_public_settings(include_data=True)}}


@router.delete("/settings/background")
def delete_template_report_background(user=Depends(require_page_manager("templatereport"))):
    now = dt.datetime.now(dt.timezone.utc).isoformat()
    with _SETTINGS_LOCK:
        BACKGROUND_FILE.unlink(missing_ok=True)
        settings = _load_settings()
        settings.update({
            "background_updated_by": str(user.get("username") or ""),
            "background_updated_at": now,
        })
        save_json(SETTINGS_FILE, settings)
    return {"ok": True, "settings": {"background": _background_public_settings()}}


# ── 저장 ──────────────────────────────────────────────────────────────────────
def _save_slot(slot: TemplateSlotReq, page_index: int, history: dict, history_by_name: dict, old_slots: dict) -> dict:
    kind = _slot_kind(slot)
    position = int(slot.position)
    common = {
        "position": position,
        "kind": kind,
        "title": _clean_text(slot.title, 180),
    }
    if kind == "chart":
        chart_ref = _clean_text(slot.chart_id, 120)
        if not chart_ref and not str(slot.definition_code or "").strip():
            raise HTTPException(400, f"{page_index + 1}페이지 {position}번 차트를 선택해 주세요.")
        if not chart_ref:
            chart_ref = f"inline_p{page_index + 1}_{position}"
        history_row = history.get(chart_ref) or history_by_name.get(chart_ref.casefold())
        chart_id = str((history_row or {}).get("history_id") or chart_ref)
        old_slot = old_slots.get(chart_id) or old_slots.get(chart_ref) or {}
        raw_definition = str(
            slot.definition_code
            or (history_row or {}).get("definition_code")
            or old_slot.get("definition_code")
            or ""
        )
        # 코드 복붙 왕복 시 마지막 개행까지 보존한다. 제어문자만 제거한다.
        definition = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]+", " ", raw_definition)[:100_000]
        if not definition.strip():
            raise HTTPException(400, f"Chart ID 또는 Name을 찾지 못했습니다: {chart_ref}")
        try:
            parse_chart_builder_definition(definition)
        except ChartBuilderDefinitionError as exc:
            raise HTTPException(400, f"{page_index + 1}페이지 {position}번 차트 생성식 오류: {exc}") from exc
        chart_name = _clean_text(
            (history_row or {}).get("name")
            or slot.chart_name
            or old_slot.get("chart_name")
            or f"Chart {position}",
            120,
        )
        layout = _slot_layout({
            "position": position, "kind": kind, "x": slot.x, "y": slot.y,
            "chart_width": slot.chart_width, "chart_height": slot.chart_height,
            "definition_code": definition,
        })
        return {
            **common,
            **{key: layout[key] for key in ("x", "y", "chart_width", "chart_height")},
            "chart_id": chart_id,
            "chart_name": chart_name,
            "chart_label": chart_name,
            "definition_code": definition[:100_000],
        }

    layout = _slot_layout({
        "position": position, "kind": kind,
        "x": slot.x, "y": slot.y, "width": slot.width, "height": slot.height,
    })
    block = {**common, **{key: layout[key] for key in ("x", "y", "width", "height")}}
    if kind == "text":
        text = _clean_multiline(slot.text, 4000)
        if not text:
            raise HTTPException(400, f"{page_index + 1}페이지 {position}번 글 블록의 내용을 입력해 주세요.")
        block["text"] = text
    elif kind == "split":
        product = _clean_text(slot.product, 120)
        lot = _clean_text(slot.lot, 120)
        if not product or not lot:
            raise HTTPException(400, f"{page_index + 1}페이지 {position}번 Split 블록에는 제품과 랏이 필요합니다.")
        block.update({
            "product": product,
            "lot": lot,
            "columns": _clean_text(slot.columns, 2000),
            "display_mode": "split_check" if str(slot.display_mode or "").strip() == "split_check" else "matrix",
        })
    elif kind in {"stats", "legend"}:
        source = int(slot.source_position or 0)
        if source < 1:
            label = "통계표" if kind == "stats" else "범례"
            raise HTTPException(400, f"{page_index + 1}페이지 {position}번 {label}는 대상 차트 번호가 필요합니다.")
        block["source_position"] = source
        if kind == "stats":
            block["stats"] = _clean_text(slot.stats, 200) or "n,mean,median,std"
    return block


@router.post("/templates")
def save_template(req: TemplateSaveReq, user=Depends(current_user)):
    requested_name = _clean_text(req.name, 120)
    if not requested_name:
        raise HTTPException(400, "Template 이름을 입력해 주세요.")
    if not req.pages or len(req.pages) > MAX_PAGES:
        raise HTTPException(400, f"페이지는 1~{MAX_PAGES}개까지 저장할 수 있습니다.")

    history = _chart_history()
    history_by_name = _chart_history_by_name(history)
    with _STORE_LOCK:
        rows = _load_templates()
        existing = next((row for row in rows if str(row.get("id")) == str(req.id)), None)
        if existing and user.get("role") != "admin" and existing.get("created_by") != user.get("username"):
            raise HTTPException(403, "작성자 또는 관리자만 이 Template을 수정할 수 있습니다.")
        old_slots = {
            str(slot.get("chart_id")): slot
            for page in ((existing or {}).get("pages") or [])
            for slot in (page.get("slots") or [])
            if isinstance(slot, dict) and slot.get("chart_id")
        }
        used_names = {
            str(value).casefold()
            for row in rows
            if not existing or str(row.get("id")) != str(existing.get("id"))
            for value in (row.get("id"), row.get("name"))
            if value
        }
        name = _unique_template_name(requested_name, used_names)
        pages = []
        for page_index, page in enumerate(req.pages):
            if len(page.slots) > MAX_CHARTS_PER_PAGE:
                raise HTTPException(400, f"페이지당 블록은 {MAX_CHARTS_PER_PAGE}개까지 배치할 수 있습니다.")
            positions: set[int] = set()
            slots = []
            for slot in page.slots:
                position = int(slot.position)
                if position < 1 or position > 1000 or position in positions:
                    raise HTTPException(400, f"{page_index + 1}페이지의 차트 위치가 올바르지 않습니다.")
                positions.add(position)
                slots.append(_save_slot(slot, page_index, history, history_by_name, old_slots))
            chart_positions = {int(slot["position"]) for slot in slots if slot["kind"] == "chart"}
            for slot in slots:
                if slot["kind"] in {"stats", "legend"} and int(slot["source_position"]) not in chart_positions:
                    label = "통계표" if slot["kind"] == "stats" else "범례"
                    raise HTTPException(400, f"{page_index + 1}페이지 {label}가 가리키는 차트 번호가 없습니다.")
            pages.append({
                "id": _clean_text(page.id, 80) or f"page_{uuid.uuid4().hex[:10]}",
                "title": _clean_text(page.title, 180) or f"Page {page_index + 1}",
                "subtitle": _clean_text(page.subtitle, 180),
                "slots": sorted(slots, key=lambda item: item["position"]),
            })
        options = (req.options or TemplateOptionsReq()).model_dump()
        variables = [
            {
                "name": normalize_name(item.name),
                "label": _clean_text(item.label, 80) or normalize_name(item.name),
                "default": _clean_text(item.default, 200),
            }
            for item in (req.variables or [])
            if normalize_name(item.name)
        ]
        now = dt.datetime.now(dt.timezone.utc).isoformat()
        if existing:
            template_id = str(existing.get("id") or "")
        else:
            used_ids = {str(row.get("id") or "") for row in rows if row.get("id")}
            template_id = ""
            while not template_id or template_id in used_ids:
                template_id = f"report_tpl_{uuid.uuid4().hex[:12]}"
        saved = {
            "id": template_id,
            "name": name,
            "pages": pages,
            "variables": variables,
            "options": options,
            "created_by": str((existing or {}).get("created_by") or user.get("username") or ""),
            "created_at": str((existing or {}).get("created_at") or now),
            "updated_by": str(user.get("username") or ""),
            "updated_at": now,
        }
        rows = [saved if str(row.get("id")) == template_id else row for row in rows]
        if not existing:
            rows.append(saved)
        _save_templates(rows)
    return {"ok": True, "template": _public_template(saved)}


@router.delete("/templates/{template_id}")
def delete_template(template_id: str, user=Depends(current_user)):
    with _STORE_LOCK:
        rows = _load_templates()
        target = next((row for row in rows if str(row.get("id")) == str(template_id)), None)
        if not target:
            raise HTTPException(404, "Template Report를 찾지 못했습니다.")
        if user.get("role") != "admin" and target.get("created_by") != user.get("username"):
            raise HTTPException(403, "작성자 또는 관리자만 삭제할 수 있습니다.")
        _save_templates([row for row in rows if str(row.get("id")) != str(template_id)])
    return {"ok": True, "id": template_id}


# ── 전체 Template 코드 ───────────────────────────────────────────────────────
def _template_request_from_code(code: str) -> TemplateSaveReq:
    raw = str(code or "").strip()
    if raw.startswith("```"):
        lines = raw.splitlines()
        if lines and lines[0].strip().startswith("```"):
            lines = lines[1:]
        if lines and lines[-1].strip() == "```":
            lines = lines[:-1]
        raw = "\n".join(lines).strip()
    if not raw:
        raise HTTPException(400, "Template 전체 코드를 입력해 주세요.")
    try:
        obj = json.loads(raw)
    except json.JSONDecodeError as exc:
        raise HTTPException(400, f"Template JSON {exc.lineno}행 {exc.colno}열 오류: {exc.msg}") from exc
    if isinstance(obj, dict) and isinstance(obj.get("template"), dict):
        obj = obj["template"]
    if not isinstance(obj, dict):
        raise HTTPException(400, "Template 코드는 JSON object여야 합니다.")
    try:
        return TemplateSaveReq(**obj)
    except Exception as exc:
        raise HTTPException(400, f"Template 코드 구조 오류: {exc}") from exc


def _normalize_code_template(req: TemplateSaveReq, user: dict | None = None) -> dict:
    name = _clean_text(req.name, 120)
    if not name:
        raise HTTPException(400, "Template 이름이 필요합니다.")
    if not req.pages or len(req.pages) > MAX_PAGES:
        raise HTTPException(400, f"페이지는 1~{MAX_PAGES}개까지 작성할 수 있습니다.")
    history = _chart_history()
    history_by_name = _chart_history_by_name(history)
    pages = []
    for page_index, page in enumerate(req.pages):
        if len(page.slots) > MAX_CHARTS_PER_PAGE:
            raise HTTPException(400, f"페이지당 블록은 {MAX_CHARTS_PER_PAGE}개까지 작성할 수 있습니다.")
        positions: set[int] = set()
        slots = []
        for slot in page.slots:
            position = int(slot.position)
            if position < 1 or position > 1000 or position in positions:
                raise HTTPException(400, f"{page_index + 1}페이지의 블록 번호가 올바르지 않습니다.")
            positions.add(position)
            slots.append(_save_slot(slot, page_index, history, history_by_name, {}))
        chart_positions = {int(slot["position"]) for slot in slots if slot["kind"] == "chart"}
        for slot in slots:
            if slot["kind"] in {"stats", "legend"} and int(slot["source_position"]) not in chart_positions:
                label = "통계표" if slot["kind"] == "stats" else "범례"
                raise HTTPException(400, f"{page_index + 1}페이지 {label}가 가리키는 차트 번호가 없습니다.")
        pages.append({
            "id": _clean_text(page.id, 80) or f"page_{page_index + 1}",
            "title": _clean_text(page.title, 180) or f"Page {page_index + 1}",
            "subtitle": _clean_text(page.subtitle, 180),
            "slots": sorted(slots, key=lambda item: item["position"]),
        })
    options = (req.options or TemplateOptionsReq()).model_dump()
    variables = [
        {
            "name": normalize_name(item.name),
            "label": _clean_text(item.label, 80) or normalize_name(item.name),
            "default": _clean_text(item.default, 200),
        }
        for item in (req.variables or [])
        if normalize_name(item.name)
    ]
    row = {
        "id": _clean_text(req.id, 120),
        "name": name,
        "pages": pages,
        "variables": variables,
        "options": options,
        "created_by": str((user or {}).get("username") or ""),
        "created_at": "",
        "updated_by": str((user or {}).get("username") or ""),
        "updated_at": "",
    }
    return _public_template(row)


@router.post("/code/parse")
def parse_template_code(req: TemplateCodeReq, user=Depends(current_user)):
    """Validate full JSON code without saving it, then return a canonical draft."""
    return {"ok": True, "template": _normalize_code_template(_template_request_from_code(req.code), user)}


@router.post("/assistant")
def template_assistant(req: TemplateAssistantReq, user=Depends(current_user)):
    """Let the configured company LLM create/edit full Template code, then validate it."""
    instruction = _clean_multiline(req.instruction, 3000)
    if not instruction:
        raise HTTPException(400, "AI에게 만들거나 수정할 Template 내용을 입력해 주세요.")
    current = _normalize_code_template(_template_request_from_code(req.template_code), user)
    try:
        from core import llm_adapter

        if not llm_adapter.is_available():
            return {
                "ok": True,
                "changed": False,
                "message": "연결된 사내 AI가 없어 현재 코드는 그대로 유지했습니다. 전체 코드는 직접 편집할 수 있습니다.",
                "template": current,
                "llm": {"available": False, "used": False},
            }
        system = """You create or edit a Flow semiconductor Template Report as one JSON object.
Return only an object with keys message and template. Preserve all unrelated fields.
The template must contain name, options, variables, and 1-30 pages. Each page contains title, subtitle, and slots.
Allowed slot kinds: chart, split, text. A chart slot must contain position, x, y, chart_width, chart_height, chart_name, and a complete valid definition_code in Flow ChartBuilder DSL. chart_id may reference an existing id or be blank for an inline chart.
Layout uses a 1920x1080 design: x/y are percentages and chart_width/chart_height are pixels. Every slot must stay inside the slide.
ChartBuilder DSL uses Q1/TABLE/PRODUCT/SQL, optional JOIN, then CHART with TYPE/X/Y/COLOR/TRELLIS/WIDTH/HEIGHT and MAX_ROWS.
Prefer readable one-page semiconductor meeting reports with bold axes and chart diversity when requested. Never invent database secrets or unsupported slot kinds."""
        available = list_charts(user).get("charts", [])[:40]
        payload = {
            "instruction": instruction,
            "current_template": current,
            "available_charts": [
                {
                    "id": item.get("id"),
                    "name": item.get("name"),
                    "label": item.get("label"),
                    "chart": item.get("chart"),
                    "definition_code": str(item.get("definition_code") or "")[:12_000],
                }
                for item in available
            ],
        }
        out = llm_adapter.complete_json(
            json.dumps(payload, ensure_ascii=False),
            system=system,
            timeout=45,
            max_retries=1,
            schema={
                "keys": ["message", "template"],
                "required": ["message", "template"],
                "properties": {"message": {}, "template": {}},
            },
        )
        obj = out.get("obj") if isinstance(out.get("obj"), dict) else {}
        candidate = obj.get("template") if isinstance(obj.get("template"), dict) else None
        if not out.get("ok") or not candidate:
            return {
                "ok": True,
                "changed": False,
                "message": str(obj.get("message") or out.get("error") or "AI가 유효한 Template 코드를 만들지 못했습니다."),
                "template": current,
                "llm": {"available": True, "used": False},
            }
        normalized = _normalize_code_template(TemplateSaveReq(**candidate), user)
        return {
            "ok": True,
            "changed": normalized != current,
            "message": _clean_text(obj.get("message") or "AI가 Template 전체 코드를 만들었습니다.", 500),
            "template": normalized,
            "llm": {"available": True, "used": True},
        }
    except HTTPException:
        raise
    except Exception as exc:
        raise HTTPException(400, f"AI Template 생성 실패: {exc}") from exc


# ── 실행(전개) ────────────────────────────────────────────────────────────────
def _default_page_subtitle(username: str) -> str:
    """빈 우측 표기에 쓸 실행 시점 라벨 (로컬 날짜 + 로그인 사용자)."""
    stamp = dt.datetime.now().astimezone().strftime("%Y-%m-%d")
    clean_username = _clean_text(username, 120)
    return f"{stamp} {clean_username}".strip()


def _run_params(template: dict, req, user: dict | None = None) -> dict:
    options = _template_options(template)
    try:
        bindings = validate_bindings(getattr(req, "bindings", None) or {})
    except ReportVariableError as exc:
        raise HTTPException(400, str(exc)) from exc
    for item in _template_variables(template):
        bindings.setdefault(item["name"], item["default"])

    repeat_variable = options["repeat_variable"]
    repeat_values = [
        value
        for raw in (getattr(req, "repeat_values", None) or [])
        for value in split_list(raw)
    ]
    if not repeat_values and repeat_variable and bindings.get(repeat_variable):
        repeat_values = split_list(bindings[repeat_variable])
    if len(repeat_values) > MAX_REPEAT_VALUES:
        raise HTTPException(400, f"한 번에 실행할 값은 {MAX_REPEAT_VALUES}개까지입니다.")

    raw_context = getattr(req, "context", None)
    if hasattr(raw_context, "model_dump"):
        raw_context = raw_context.model_dump()
    elif hasattr(raw_context, "dict"):
        raw_context = raw_context.dict()
    raw_context = raw_context if isinstance(raw_context, dict) else {}

    def _context_values(name: str) -> list[str]:
        values: list[str] = []
        seen: set[str] = set()
        for raw in raw_context.get(name) or []:
            value = _clean_text(raw, 160)
            key = value.casefold()
            if value and key not in seen:
                seen.add(key)
                values.append(value)
            if len(values) >= 200:
                break
        return values

    recent_days = max(0, min(3650, int(raw_context.get("recent_days") or 0)))
    color_rules = [_clean_text(rule, 500) for rule in (raw_context.get("color_rules") or []) if _clean_text(rule, 500)][:300]
    context = {
        "root_lot_ids": _context_values("root_lot_ids"),
        "wafer_ids": _context_values("wafer_ids"),
        "override_recent_days": bool(raw_context.get("override_recent_days")),
        "recent_days": recent_days,
        "date_column": _clean_text(raw_context.get("date_column") or "tkout_time", 120) or "tkout_time",
        "color_rules": color_rules,
        "color_else": _clean_text(raw_context.get("color_else") or "gray", 80) or "gray",
    }
    username = _clean_text(
        (user or {}).get("username")
        or template.get("updated_by")
        or template.get("created_by")
        or "",
        120,
    )
    return {
        "bindings": bindings,
        "repeat_variable": repeat_variable,
        "repeat_values": repeat_values or [""],
        "options": options,
        "context": context,
        "default_page_subtitle": _default_page_subtitle(username),
    }


def _resolve(text: str, bindings: dict, where: str) -> str:
    try:
        return substitute(text, bindings)
    except ReportVariableError as exc:
        raise HTTPException(400, f"{where}: {exc}") from exc


def _chart_request(definition: str, where: str, context: dict | None = None) -> dict:
    """저장 코드를 실행 요청으로 옮기고 선택된 공통 실행 컨텍스트만 덧씌운다."""
    try:
        parsed = parse_chart_builder_definition(definition)
    except ChartBuilderDefinitionError as exc:
        raise HTTPException(400, f"{where} Chart 코드 오류: {exc}") from exc
    context = context or {}
    sources = [dict(source) for source in (parsed.get("sources") or [])]
    for source in sources:
        if context.get("root_lot_ids"):
            source["runtime_root_lot_ids"] = list(context["root_lot_ids"])
        if context.get("wafer_ids"):
            source["runtime_wafer_ids"] = list(context["wafer_ids"])
        if context.get("override_recent_days"):
            days = int(context.get("recent_days") or 0)
            source["runtime_recent_days"] = days
            source["runtime_date_column"] = str(context.get("date_column") or "tkout_time") if days else ""
    chart = dict(parsed.get("chart") or {})
    if context.get("color_rules"):
        chart.update({
            "color": "custom",
            "color_rules": list(context["color_rules"]),
            "color_else": str(context.get("color_else") or "gray"),
        })
    context_pairs = linked_chart_color_pairs(chart) if context.get("color_rules") else []
    for source in sources:
        if context_pairs:
            source["runtime_lot_wafer_pairs"] = [dict(pair) for pair in context_pairs]
            if not context.get("root_lot_ids"):
                source["runtime_root_lot_ids"] = list(dict.fromkeys(pair["root_lot_id"] for pair in context_pairs))
            if not context.get("wafer_ids"):
                source["runtime_wafer_ids"] = list(dict.fromkeys(pair["wafer_id"] for pair in context_pairs))
        elif context.get("root_lot_ids") or context.get("wafer_ids"):
            # An explicit report scope replaces any exact pairs saved in the chart.
            source["runtime_lot_wafer_pairs"] = []
    return {
        "sources": sources,
        "joins": parsed.get("joins") or [],
        "max_rows": parsed.get("max_rows") or 10000,
        "chart": chart,
        "save_history": False,
    }


def _expand_deck(template: dict, params: dict) -> dict:
    """템플릿 + 실행 인자 → 실제로 그릴 페이지 목록. 인자가 같으면 결과도 같다.

    미리보기와 PPTX 내보내기가 같은 함수를 다시 부르기 때문에, 화면에서 본 자리와
    파일의 자리가 어긋나지 않는다.
    """
    options = params["options"]
    pages_out: list[dict] = []
    charts: list[dict] = []
    for repeat_value in params["repeat_values"]:
        bindings = dict(params["bindings"])
        if params["repeat_variable"] and repeat_value:
            bindings[params["repeat_variable"]] = repeat_value
        for page in template.get("pages") or []:
            page_index = len(pages_out)
            if page_index >= MAX_RENDERED_PAGES:
                raise HTTPException(400, f"생성할 페이지가 {MAX_RENDERED_PAGES}장을 넘습니다. 값 개수를 줄여 주세요.")
            where = f"{page_index + 1}페이지"
            blocks: list[dict] = []
            for slot in page.get("slots") or []:
                position = int(slot.get("position") or 1)
                kind = _slot_kind(slot)
                layout = _slot_layout(slot)
                block = {
                    "key": f"{page_index}:{position}",
                    "page_index": page_index,
                    "position": position,
                    "kind": kind,
                    "title": _resolve(str(slot.get("title") or ""), bindings, where),
                    "repeat_value": repeat_value,
                    **{name: layout[name] for name in ("x", "y", "width", "height", "chart_width", "chart_height")},
                }
                if kind == "chart":
                    definition = _resolve(str(slot.get("definition_code") or ""), bindings, where)
                    block.update({
                        "chart_id": str(slot.get("chart_id") or ""),
                        "chart_label": str(slot.get("chart_label") or slot.get("chart_name") or ""),
                        "request": _chart_request(definition, where, params.get("context")),
                    })
                    charts.append(block)
                elif kind == "text":
                    block["text"] = _resolve(str(slot.get("text") or ""), bindings, where)
                elif kind == "split":
                    block.update({
                        "product": _resolve(str(slot.get("product") or ""), bindings, where),
                        "lot": _resolve(str(slot.get("lot") or ""), bindings, where),
                        "columns": _resolve(str(slot.get("columns") or ""), bindings, where),
                        "display_mode": str(slot.get("display_mode") or "matrix"),
                    })
                elif kind in {"stats", "legend"}:
                    source = int(slot.get("source_position") or 0)
                    block.update({
                        "source_position": source,
                        "source_key": f"{page_index}:{source}",
                    })
                    if kind == "stats":
                        block["stats"] = str(slot.get("stats") or "n,mean,median,std")
                blocks.append(block)
            page_subtitle = _resolve(str(page.get("subtitle") or ""), bindings, where)
            pages_out.append({
                "index": page_index,
                "title": _resolve(str(page.get("title") or f"Page {page_index + 1}"), bindings, where),
                "subtitle": page_subtitle if page_subtitle.strip() else params["default_page_subtitle"],
                "repeat_value": repeat_value,
                "blocks": blocks,
            })
    cover_bindings = dict(params["bindings"])
    if params["repeat_variable"] and params["repeat_values"][0]:
        cover_bindings[params["repeat_variable"]] = params["repeat_values"][0]
    subtitle = _resolve(options["subtitle"], cover_bindings, "표지") if options["subtitle"] else ""
    if not subtitle:
        parts = [value for value in params["repeat_values"] if value]
        subtitle = " · ".join(parts[:4]) + ("…" if len(parts) > 4 else "")
    return {
        "title": str(template.get("name") or "Template Report"),
        "subtitle": subtitle,
        "cover": options["cover"],
        "footer": options["footer"],
        "pages": pages_out,
        "charts": charts,
        "bindings": dict(params.get("bindings") or {}),
        "repeat_variable": str(params.get("repeat_variable") or ""),
        "repeat_values": [str(value) for value in (params.get("repeat_values") or []) if str(value)],
        "context": dict(params.get("context") or {}),
    }


@router.post("/run")
def prepare_run(req: TemplateRunReq, user=Depends(current_user)):
    template = _template_or_404(req.template_id)
    params = _run_params(template, req, user)
    deck = _expand_deck(template, params)
    return {
        "ok": True,
        "template": _public_template(template),
        "bindings": params["bindings"],
        "repeat_variable": params["repeat_variable"],
        "repeat_values": [value for value in params["repeat_values"] if value],
        "context": params["context"],
        "deck": deck,
        "charts": deck["charts"],
    }


@router.post("/split-table")
def split_table_block(req: SplitBlockReq, _user=Depends(current_user)):
    """Split 블록 내용 — Inform 스냅샷과 같은 SplitTable 화면 규약을 그대로 쓴다."""
    from app_v2.modules.informs.splittable_embed import build_splittable_embed

    columns = [item.strip() for item in str(req.columns or "").split(",") if item.strip()]
    embed = build_splittable_embed(product=req.product, lot_id=req.lot_id, custom_cols=columns)
    header = ["parameter", *(embed.get("st_view", {}).get("headers") or [])]
    rows = [[str(cell) for cell in row] for row in (embed.get("rows") or [])]
    max_rows = max(1, min(MAX_TABLE_ROWS * 4, int(req.max_rows or MAX_TABLE_ROWS)))
    truncated = len(rows) > max_rows
    return {
        "ok": True,
        "columns": header,
        "rows": rows[:max_rows],
        "row_total": len(rows),
        "truncated": truncated,
        "note": str(embed.get("note") or ""),
        "source": str(embed.get("source") or ""),
    }


# ── 내보내기 ──────────────────────────────────────────────────────────────────
def _image_key(image: ExportImageReq) -> str:
    key = str(image.key or "").strip()
    return key or f"{int(image.page_index)}:{int(image.position)}"


def _decode_images(images: list[ExportImageReq]) -> dict[str, bytes]:
    if len(images) > MAX_IMAGES:
        raise HTTPException(400, "차트 이미지가 허용 개수를 넘었습니다.")
    decoded: dict[str, bytes] = {}
    for image in images:
        match = re.fullmatch(r"data:image/(png|jpeg);base64,([A-Za-z0-9+/=\s]+)", str(image.data_url or ""), re.IGNORECASE)
        if not match:
            raise HTTPException(400, "차트 이미지는 PNG 또는 JPEG data URL이어야 합니다.")
        try:
            payload = base64.b64decode(match.group(2), validate=False)
        except Exception as exc:
            raise HTTPException(400, "차트 이미지 base64를 해석하지 못했습니다.") from exc
        if not payload or len(payload) > 12 * 1024 * 1024:
            raise HTTPException(400, "차트 이미지 한 장은 12MB 이하여야 합니다.")
        decoded[_image_key(image)] = payload
    return decoded


def _decode_tables(tables: list[ExportTableReq]) -> dict[str, dict]:
    out: dict[str, dict] = {}
    for table in tables or []:
        key = str(table.key or "").strip()
        if not key:
            continue
        columns = [_clean_text(value, 60) for value in (table.columns or [])][:MAX_TABLE_COLUMNS]
        rows = [
            [_clean_text(cell, 60) for cell in (row or [])][:MAX_TABLE_COLUMNS]
            for row in (table.rows or [])[:MAX_TABLE_ROWS]
        ]
        out[key] = {
            "title": _clean_text(table.title, 180),
            "columns": columns,
            "rows": rows,
            "note": _clean_text(table.note, 200),
        }
    return out


def _download_header(filename: str) -> str:
    """Return an ASCII-safe Content-Disposition with a UTF-8 filename fallback."""
    cleaned = safe_filename(filename)
    ascii_name = cleaned.encode("ascii", "ignore").decode("ascii").strip("._") or "download"
    suffix = Path(cleaned).suffix
    if suffix and not ascii_name.lower().endswith(suffix.lower()):
        ascii_name += suffix
    return f"attachment; filename=\"{ascii_name}\"; filename*=UTF-8''{quote(cleaned)}"


def _rgb(value):
    from pptx.dml.color import RGBColor

    return RGBColor(*value)


def _accent():
    return _rgb(REPORT_ACCENT)


def _text_color():
    return _rgb(REPORT_TEXT)


def _muted_color():
    return _rgb(REPORT_MUTED)


def _navy():
    """Legacy helper kept for older callers; Flow's brand accent is now orange."""
    return _accent()


def _add_bar(slide, top_in: float, height_in: float, width_emu):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(top_in), width_emu, Inches(height_in))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _accent()
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


def _write(text_frame, text: str, *, size: int, bold: bool, color, align, font: str = REPORT_FONT):
    from pptx.util import Pt

    text_frame.clear()
    lines = str(text or "").split("\n")
    for index, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = line
        paragraph.alignment = align
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.name = font
        paragraph.font.color.rgb = color


def _add_title_bar(slide, deck_width_emu, title: str, subtitle: str):
    """내부 회의용 얇은 헤더. 제목보다 차트와 축의 가독성을 우선한다."""
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    title_box = slide.shapes.add_textbox(Inches(0.24), Inches(0.03), Inches(9.5), Inches(TITLE_BAR_HEIGHT_IN - 0.07))
    title_box.text_frame.word_wrap = False
    title_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    _write(title_box.text_frame, title, size=14, bold=True, color=_text_color(), align=PP_ALIGN.LEFT)
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(9.8), Inches(0.04), Inches(3.25), Inches(TITLE_BAR_HEIGHT_IN - 0.08))
        sub_box.text_frame.word_wrap = False
        sub_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        _write(sub_box.text_frame, subtitle, size=9, bold=False, color=_muted_color(), align=PP_ALIGN.RIGHT)
    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.22), Inches(TITLE_BAR_HEIGHT_IN - 0.015), Inches(SLIDE_WIDTH_IN - 0.44), Pt(0.8))
    divider.fill.solid()
    divider.fill.fore_color.rgb = _accent()
    divider.line.fill.background()


def _add_footer(slide, left_text: str, right_text: str):
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches

    grey = _muted_color()
    if str(left_text or "").strip():
        left_box = slide.shapes.add_textbox(Inches(0.3), Inches(SLIDE_HEIGHT_IN - 0.34), Inches(8.5), Inches(0.26))
        left_box.text_frame.word_wrap = False
        _write(left_box.text_frame, left_text, size=9, bold=False, color=grey, align=PP_ALIGN.LEFT)
    right_box = slide.shapes.add_textbox(Inches(9.5), Inches(SLIDE_HEIGHT_IN - 0.34), Inches(3.5), Inches(0.26))
    right_box.text_frame.word_wrap = False
    _write(right_box.text_frame, right_text, size=9, bold=False, color=grey, align=PP_ALIGN.RIGHT)


def _add_background_picture(slide, deck, payload: bytes):
    """배경을 16:9 슬라이드에 cover 방식으로 깐다(찌그러뜨리지 않고 중앙 crop)."""
    if not payload:
        return None
    try:
        from PIL import Image

        with Image.open(io.BytesIO(payload)) as image:
            image_width, image_height = image.size
    except (OSError, ValueError):
        return None
    if image_width < 1 or image_height < 1:
        return None
    slide_width, slide_height = int(deck.slide_width), int(deck.slide_height)
    image_ratio = image_width / image_height
    slide_ratio = slide_width / slide_height
    if image_ratio >= slide_ratio:
        height = slide_height
        width = round(height * image_ratio)
        left, top = (slide_width - width) // 2, 0
    else:
        width = slide_width
        height = round(width / image_ratio)
        left, top = 0, (slide_height - height) // 2
    return slide.shapes.add_picture(io.BytesIO(payload), left, top, width=width, height=height)


def _add_cover(deck, title: str, subtitle: str, meta: str, *, background_image: bytes = b""):
    """표지 — auto report `make_title_page` 와 같은 구성(상하 바·대제목·구분선)."""
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    slide = deck.slides.add_slide(deck.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(REPORT_PAGE)
    _add_background_picture(slide, deck, background_image)
    grey = _muted_color()
    _add_bar(slide, 0.0, COVER_BAR_HEIGHT_IN, deck.slide_width)
    _add_bar(slide, SLIDE_HEIGHT_IN - COVER_BAR_HEIGHT_IN, COVER_BAR_HEIGHT_IN, deck.slide_width)

    date_box = slide.shapes.add_textbox(Inches(SLIDE_WIDTH_IN - 3.9), Inches(0.7), Inches(3.6), Inches(0.4))
    _write(date_box.text_frame, dt.datetime.now().strftime("%Y-%m-%d"), size=13, bold=False, color=grey, align=PP_ALIGN.RIGHT)

    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(2.55), Inches(SLIDE_WIDTH_IN - 1.2), Inches(1.5))
    title_box.text_frame.word_wrap = True
    title_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    _write(title_box.text_frame, title, size=44, bold=True, color=_text_color(), align=PP_ALIGN.CENTER)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(SLIDE_WIDTH_IN / 2 - 2.2), Inches(4.18), Inches(4.4), Pt(2.5))
    line.fill.solid()
    line.fill.fore_color.rgb = _accent()
    line.line.fill.background()
    line.shadow.inherit = False

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(4.35), Inches(SLIDE_WIDTH_IN - 1.2), Inches(0.8))
        sub_box.text_frame.word_wrap = True
        _write(sub_box.text_frame, subtitle, size=22, bold=False, color=grey, align=PP_ALIGN.CENTER)
    if meta:
        meta_box = slide.shapes.add_textbox(Inches(0.6), Inches(5.25), Inches(SLIDE_WIDTH_IN - 1.2), Inches(0.5))
        meta_box.text_frame.word_wrap = True
        _write(meta_box.text_frame, meta, size=13, bold=False, color=grey, align=PP_ALIGN.CENTER)
    return slide


def _set_cell_border(cell):
    """Carbon data table처럼 세로선 없이 아래 구분선만 둔다."""
    from pptx.oxml.xmlchemy import OxmlElement

    tcPr = cell._tc.get_or_add_tcPr()
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    for tag in ("lnL", "lnR", "lnT", "lnB"):
        for element in tcPr.findall(ns + tag):
            tcPr.remove(element)
    line = OxmlElement("a:lnB")
    line.set("w", "6350")
    line.set("cmpd", "sng")
    fill = OxmlElement("a:solidFill")
    color = OxmlElement("a:srgbClr")
    color.set("val", "E5E5E5")
    fill.append(color)
    line.append(fill)
    tcPr.insert(0, line)


def _add_table(slide, rect: dict, table: dict):
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    columns = list(table.get("columns") or [])
    rows = [list(row) for row in (table.get("rows") or [])]
    if not columns and rows:
        columns = [f"C{index + 1}" for index in range(len(rows[0]))]
    if not columns:
        return
    columns = columns[:MAX_TABLE_COLUMNS]
    rows = [row[:len(columns)] + [""] * max(0, len(columns) - len(row)) for row in rows[:MAX_TABLE_ROWS]]

    x = SLIDE_WIDTH_IN * float(rect["x"]) / 100
    y = SLIDE_HEIGHT_IN * float(rect["y"]) / 100
    width = SLIDE_WIDTH_IN * float(rect["width"]) / 100
    height = SLIDE_HEIGHT_IN * float(rect["height"]) / 100
    bottom = y + height
    title = str(table.get("title") or "")
    if title:
        label = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(0.26))
        label.text_frame.word_wrap = False
        _write(label.text_frame, title, size=12, bold=True, color=_text_color(), align=PP_ALIGN.LEFT)
        y += 0.3
        height = max(0.4, height - 0.3)

    note = str(table.get("note") or "")
    if note:
        height = max(0.4, height - 0.24)
    count = len(rows) + 1
    row_height = max(0.15, min(0.3, height / count))
    shape = slide.shapes.add_table(count, len(columns), Inches(x), Inches(y), Inches(width), Inches(row_height * count))
    grid = shape.table
    # 첫 열은 항목 이름이라 넓게, 나머지는 균등 — SplitTable 화면과 같은 감각.
    first = min(width * 0.34, 2.6) if len(columns) > 2 else width / len(columns)
    rest = (width - first) / max(1, len(columns) - 1) if len(columns) > 1 else width
    grid.columns[0].width = Inches(first)
    for index in range(1, len(columns)):
        grid.columns[index].width = Inches(rest)
    for index in range(count):
        grid.rows[index].height = Inches(row_height)

    font_size = 9 if len(columns) <= 8 else (8 if len(columns) <= 14 else 7)
    for index, name in enumerate(columns):
        cell = grid.cell(0, index)
        cell.text = str(name)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(REPORT_SUBTLE)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = cell.margin_right = Inches(0.02)
        cell.margin_top = cell.margin_bottom = Inches(0.0)
        _set_cell_border(cell)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(font_size)
            paragraph.font.bold = True
            paragraph.font.name = REPORT_FONT
            paragraph.font.color.rgb = _text_color()
            paragraph.alignment = PP_ALIGN.CENTER
    for row_index, row in enumerate(rows, start=1):
        stripe = _rgb(REPORT_PAGE) if row_index % 2 == 0 else _rgb(REPORT_PANEL)
        for column_index, value in enumerate(row):
            cell = grid.cell(row_index, column_index)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = stripe
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = cell.margin_right = Inches(0.02)
            cell.margin_top = cell.margin_bottom = Inches(0.0)
            cell.text_frame.word_wrap = False
            _set_cell_border(cell)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.bold = column_index == 0
                paragraph.font.name = REPORT_FONT
                paragraph.font.color.rgb = _text_color()
                paragraph.alignment = PP_ALIGN.LEFT if column_index == 0 else PP_ALIGN.CENTER
    if note:
        # PowerPoint 는 글자가 길면 행을 스스로 늘린다. 머리글이 두 줄이 되는 경우까지
        # 감안해 한 줄치 여유를 두되, 블록 바깥으로는 내려가지 않게 잡는다.
        note_top = min(y + row_height * (count + 0.8) + 0.02, max(y, bottom - 0.24))
        note_box = slide.shapes.add_textbox(Inches(x), Inches(note_top), Inches(width), Inches(0.22))
        note_box.text_frame.word_wrap = False
        _write(note_box.text_frame, note, size=9, bold=False, color=RGBColor(0x70, 0x7A, 0x8A), align=PP_ALIGN.LEFT)


def _add_text_block(slide, rect: dict, title: str, text: str):
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches

    x = SLIDE_WIDTH_IN * float(rect["x"]) / 100
    y = SLIDE_HEIGHT_IN * float(rect["y"]) / 100
    width = SLIDE_WIDTH_IN * float(rect["width"]) / 100
    height = SLIDE_HEIGHT_IN * float(rect["height"]) / 100
    if title:
        label = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(0.26))
        label.text_frame.word_wrap = False
        _write(label.text_frame, title, size=12, bold=True, color=_text_color(), align=PP_ALIGN.LEFT)
        y += 0.3
        height = max(0.3, height - 0.3)
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    box.text_frame.word_wrap = True
    _write(box.text_frame, text, size=12, bold=False, color=RGBColor(0x1F, 0x29, 0x37), align=PP_ALIGN.LEFT)


def _ppt_rgb(value: str, fallback=(99, 102, 241)):
    from pptx.dml.color import RGBColor

    raw = str(value or "").strip().casefold()
    named = {
        "red": (239, 68, 68), "blue": (59, 130, 246), "green": (16, 185, 129),
        "orange": (245, 158, 11), "yellow": (234, 179, 8), "purple": (139, 92, 246),
        "pink": (236, 72, 153), "gray": (148, 163, 184), "grey": (148, 163, 184),
        "black": (15, 23, 42), "white": (255, 255, 255),
    }
    if raw in named:
        return RGBColor(*named[raw])
    match = re.fullmatch(r"#?([0-9a-f]{6})", raw)
    if match:
        token = match.group(1)
        return RGBColor(int(token[0:2], 16), int(token[2:4], 16), int(token[4:6], 16))
    return RGBColor(*fallback)


def _add_legend(slide, rect: dict, table: dict):
    """공통 범례를 PPT 도형으로 그린다. PNG보다 선명하고 파일도 작다."""
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches

    x = SLIDE_WIDTH_IN * float(rect["x"]) / 100
    y = SLIDE_HEIGHT_IN * float(rect["y"]) / 100
    width = SLIDE_WIDTH_IN * float(rect["width"]) / 100
    height = SLIDE_HEIGHT_IN * float(rect["height"]) / 100
    title = str(table.get("title") or _block_label(rect) or "Legend")
    if title:
        title_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(0.24))
        title_box.text_frame.word_wrap = False
        _write(title_box.text_frame, title, size=11, bold=True, color=_text_color(), align=PP_ALIGN.LEFT)
        y += 0.28
        height = max(0.25, height - 0.28)

    rows = [list(row) for row in (table.get("rows") or [])][:24]
    if not rows:
        return
    columns = 1 if width < 2.2 else (2 if width < 4.4 else 3)
    cell_width = width / columns
    row_count = math.ceil(len(rows) / columns)
    row_height = min(0.28, height / max(1, row_count))
    for index, row in enumerate(rows):
        column = index % columns
        line = index // columns
        item_x = x + column * cell_width
        item_y = y + line * row_height
        label = str(row[0] if row else "")
        color = str(row[1] if len(row) > 1 else "")
        count = str(row[2] if len(row) > 2 else "").strip()
        swatch = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(item_x), Inches(item_y + 0.045), Inches(0.14), Inches(0.14)
        )
        swatch.fill.solid()
        swatch.fill.fore_color.rgb = _ppt_rgb(color)
        swatch.line.fill.background()
        text_box = slide.shapes.add_textbox(
            Inches(item_x + 0.19), Inches(item_y), Inches(max(0.2, cell_width - 0.21)), Inches(row_height)
        )
        text_box.text_frame.word_wrap = False
        text_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        shown = f"{label} ({count})" if count else label
        _write(text_box.text_frame, shown, size=9, bold=False, color=_text_color(), align=PP_ALIGN.LEFT)


def _format_report_time(value: str) -> str:
    raw = str(value or "").strip()
    if not raw:
        return "—"
    try:
        parsed = dt.datetime.fromisoformat(raw.replace("Z", "+00:00"))
        if parsed.tzinfo is not None:
            parsed = parsed.astimezone()
        return parsed.strftime("%Y-%m-%d %H:%M")
    except ValueError:
        return raw[:40]


def _short_json(value, limit: int = 260) -> str:
    if value in (None, "", [], {}):
        return "—"
    rendered = json.dumps(value, ensure_ascii=False, separators=(", ", ": "))
    return rendered if len(rendered) <= limit else rendered[:limit - 1] + "…"


def _context_summary(value: dict | None) -> str:
    context = value if isinstance(value, dict) else {}
    lots = context.get("root_lot_ids") or []
    wafers = context.get("wafer_ids") or []
    if context.get("override_recent_days"):
        period = f"{int(context.get('recent_days') or 0)}d / {context.get('date_column') or 'tkout_time'}"
    else:
        period = "saved chart window"
    colors = len(context.get("color_rules") or [])
    return f"lots {_short_json(lots, 74)} · wafers {_short_json(wafers, 74)} · period {period} · color rules {colors}"


def _appendix_chart_rows(template: dict) -> list[list[str]]:
    rows: list[list[str]] = []
    for page_index, page in enumerate(template.get("pages") or [], start=1):
        for slot in page.get("slots") or []:
            if _slot_kind(slot) != "chart":
                continue
            layout = _slot_layout(slot)
            chart = {}
            try:
                chart = dict(parse_chart_builder_definition(str(slot.get("definition_code") or "")).get("chart") or {})
            except ChartBuilderDefinitionError:
                pass
            chart_bits = [str(chart.get("type") or "chart")]
            if chart.get("x") or chart.get("y"):
                chart_bits.append(f"{chart.get('x') or '—'} → {chart.get('y') or '—'}")
            if chart.get("color"):
                chart_bits.append(f"color: {chart['color']}")
            if chart.get("trellis"):
                chart_bits.append(f"trellis: {chart['trellis']}")
            chart_bits.append("legend: off" if chart.get("show_legend") is False else "legend: on")
            rows.append([
                f"{page_index} · {int(slot.get('position') or 1)}",
                str(slot.get("chart_name") or slot.get("chart_label") or slot.get("chart_id") or "Chart")[:70],
                str(slot.get("chart_id") or "inline")[:80],
                " · ".join(chart_bits)[:150] +
                f" | x {float(layout['x']):.1f}%, y {float(layout['y']):.1f}%, {int(layout['chart_width'])}×{int(layout['chart_height'])}px",
            ])
    return rows


def _add_appendix_chart_table(slide, rows: list[list[str]], *, top: float):
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    columns = ["Page · Pos", "Chart Name", "Chart ID", "Chart / Layout"]
    left, width = 0.36, SLIDE_WIDTH_IN - 0.72
    row_height = 0.24
    shape = slide.shapes.add_table(
        len(rows) + 1, len(columns), Inches(left), Inches(top), Inches(width), Inches(row_height * (len(rows) + 1))
    )
    grid = shape.table
    widths = (0.10, 0.22, 0.25, 0.43)
    for index, ratio in enumerate(widths):
        grid.columns[index].width = Inches(width * ratio)
    for row in grid.rows:
        row.height = Inches(row_height)
    for column_index, name in enumerate(columns):
        cell = grid.cell(0, column_index)
        cell.text = name
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(REPORT_SUBTLE)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        _set_cell_border(cell)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(8)
            paragraph.font.bold = True
            paragraph.font.name = REPORT_FONT
            paragraph.font.color.rgb = _text_color()
            paragraph.alignment = PP_ALIGN.LEFT
    for row_index, row in enumerate(rows, start=1):
        stripe = _rgb(REPORT_PAGE) if row_index % 2 == 0 else _rgb(REPORT_PANEL)
        for column_index, value in enumerate(row):
            cell = grid.cell(row_index, column_index)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = stripe
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = cell.margin_right = Inches(0.04)
            cell.margin_top = cell.margin_bottom = Inches(0.0)
            cell.text_frame.word_wrap = False
            _set_cell_border(cell)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(8)
                paragraph.font.bold = column_index == 0
                paragraph.font.name = REPORT_FONT
                paragraph.font.color.rgb = _text_color()
                paragraph.alignment = PP_ALIGN.LEFT


def _add_info_card(slide, *, x: float, y: float, width: float, title: str, lines: list[str]):
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches

    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(1.28))
    card.fill.solid()
    card.fill.fore_color.rgb = _rgb(REPORT_PANEL)
    card.line.color.rgb = _rgb(REPORT_BORDER)
    card.shadow.inherit = False
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.045), Inches(1.28))
    accent.fill.solid()
    accent.fill.fore_color.rgb = _accent()
    accent.line.fill.background()
    title_box = slide.shapes.add_textbox(Inches(x + 0.14), Inches(y + 0.07), Inches(width - 0.25), Inches(0.20))
    _write(title_box.text_frame, title, size=10, bold=True, color=_text_color(), align=PP_ALIGN.LEFT)
    body = slide.shapes.add_textbox(Inches(x + 0.14), Inches(y + 0.31), Inches(width - 0.25), Inches(0.88))
    body.text_frame.word_wrap = True
    body.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    _write(body.text_frame, "\n".join(lines), size=8, bold=False, color=_text_color(), align=PP_ALIGN.LEFT)


def _add_template_info_appendix(
    deck, template: dict, deck_spec: dict, *, report_user: str,
    report_generated_at: str, background_image: bytes = b"",
):
    """항상 보고서 맨 뒤에 재생성 정보만 담은 별도 Appendix를 붙인다."""
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches

    chart_rows = _appendix_chart_rows(template)
    first_count, continuation_count = 18, 24
    chunks = [chart_rows[:first_count]]
    remaining = chart_rows[first_count:]
    while remaining:
        chunks.append(remaining[:continuation_count])
        remaining = remaining[continuation_count:]
    total = len(chunks)
    blank = deck.slide_layouts[6]
    for index, rows in enumerate(chunks):
        slide = deck.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = _rgb(REPORT_PAGE)
        _add_background_picture(slide, deck, background_image)
        heading = "Template Report Info" if total == 1 else f"Template Report Info ({index + 1}/{total})"
        _add_title_bar(slide, deck.slide_width, heading, "Appendix · Reproduction metadata")
        if index == 0:
            options = template.get("options") or {}
            _add_info_card(slide, x=0.36, y=0.57, width=6.22, title="Template", lines=[
                f"Name / ID  {template.get('name') or 'Template Report'} / {template.get('id') or '—'}",
                f"Created  {template.get('created_by') or '—'} · {_format_report_time(template.get('created_at'))}",
                f"Updated  {template.get('updated_by') or '—'} · {_format_report_time(template.get('updated_at'))}",
                f"Pages / options  {len(template.get('pages') or [])} · cover {bool(options.get('cover', True))} · footer {bool(options.get('footer', True))}",
            ])
            _add_info_card(slide, x=6.74, y=0.57, width=6.22, title="Report / Recreate", lines=[
                f"Generated  {report_user or '—'} · {_format_report_time(report_generated_at)}",
                f"Bindings  {_short_json(deck_spec.get('bindings'))}",
                f"Repeat  {deck_spec.get('repeat_variable') or '—'} = {_short_json(deck_spec.get('repeat_values'))}",
                f"Context  {_context_summary(deck_spec.get('context'))}",
                f"Recreate  Open Template ID above → apply inputs → Run → export PPTX",
            ])
            table_top = 1.98
        else:
            table_top = 0.62
        label = slide.shapes.add_textbox(Inches(0.36), Inches(table_top), Inches(12.61), Inches(0.23))
        _write(label.text_frame, "Included charts · use the Template ID and Chart IDs below to reproduce this report", size=9, bold=True, color=_text_color(), align=PP_ALIGN.LEFT)
        if rows:
            _add_appendix_chart_table(slide, rows, top=table_top + 0.25)
        else:
            empty = slide.shapes.add_textbox(Inches(0.36), Inches(table_top + 0.28), Inches(12.61), Inches(0.5))
            _write(empty.text_frame, "No chart blocks are included in this template.", size=9, bold=False, color=_muted_color(), align=PP_ALIGN.LEFT)
        _add_footer(slide, "", f"Appendix {index + 1} / {total}")


def _block_label(block: dict) -> str:
    return str(block.get("title") or "").strip()


def _pptx_bytes(template: dict, images: dict[str, bytes], *,
                deck_spec: dict | None = None, tables: dict[str, dict] | None = None,
                report_user: str = "", report_generated_at: str = "",
                background_image: bytes = b"") -> bytes:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches

    tables = tables or {}
    if deck_spec is None:
        params = _run_params(template, TemplateRunReq(template_id=str(template.get("id") or "")))
        deck_spec = _expand_deck(template, params)

    deck = Presentation()
    deck.slide_width = Inches(SLIDE_WIDTH_IN)
    deck.slide_height = Inches(SLIDE_HEIGHT_IN)
    deck.core_properties.title = str(template.get("name") or "Template Report")
    deck.core_properties.subject = "Flow Template Report"
    deck.core_properties.author = str(template.get("updated_by") or template.get("created_by") or "Flow")
    blank = deck.slide_layouts[6]

    if deck_spec.get("cover"):
        _add_cover(
            deck,
            str(deck_spec.get("title") or template.get("name") or "Template Report"),
            str(deck_spec.get("subtitle") or ""),
            "Flow Template Report",
            background_image=background_image,
        )

    pages = deck_spec.get("pages") or []
    total = len(pages)
    for page in pages:
        slide = deck.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = _rgb(REPORT_PAGE)
        _add_background_picture(slide, deck, background_image)
        for block in page.get("blocks") or []:
            key = str(block.get("key") or "")
            kind = str(block.get("kind") or "chart")
            x = SLIDE_WIDTH_IN * float(block["x"]) / 100
            y = SLIDE_HEIGHT_IN * float(block["y"]) / 100
            width = SLIDE_WIDTH_IN * float(block["width"]) / 100
            height = SLIDE_HEIGHT_IN * float(block["height"]) / 100
            if kind == "chart":
                payload = images.get(key)
                if not payload:
                    continue
                slide.shapes.add_picture(io.BytesIO(payload), Inches(x), Inches(y), width=Inches(width), height=Inches(height))
            elif kind == "text":
                _add_text_block(slide, block, _block_label(block), str(block.get("text") or ""))
            else:
                table = tables.get(key)
                if not table:
                    continue
                rendered = {**table, "title": table.get("title") or _block_label(block)}
                if kind == "legend":
                    _add_legend(slide, block, rendered)
                else:
                    _add_table(slide, block, rendered)
        # TITLE stays above freely placed charts, matching the browser editor's z-order.
        _add_title_bar(slide, deck.slide_width, str(page.get("title") or ""), str(page.get("subtitle") or ""))
        if deck_spec.get("footer"):
            _add_footer(
                slide,
                "",
                f"{page.get('index', 0) + 1} / {total}",
            )
    _add_template_info_appendix(
        deck,
        template,
        deck_spec,
        report_user=report_user or str(template.get("updated_by") or template.get("created_by") or ""),
        report_generated_at=report_generated_at or dt.datetime.now(dt.timezone.utc).isoformat(),
        background_image=background_image,
    )
    out = io.BytesIO()
    deck.save(out)
    return out.getvalue()


@router.post("/export/pptx")
def export_pptx(req: ExportReq, user=Depends(current_user)):
    template = _template_or_404(req.template_id)
    params = _run_params(template, req, user)
    deck_spec = _expand_deck(template, params)
    payload = _pptx_bytes(
        template,
        _decode_images(req.images),
        deck_spec=deck_spec,
        tables=_decode_tables(req.tables),
        report_user=str(user.get("username") or ""),
        report_generated_at=dt.datetime.now(dt.timezone.utc).isoformat(),
        background_image=_background_bytes(),
    )
    stamp = dt.datetime.now().strftime("%Y%m%d")
    filename = safe_filename(f"{template.get('name') or 'template_report'}_{stamp}.pptx")
    return Response(
        content=payload,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": _download_header(filename)},
    )


@router.post("/export/images")
def export_images(req: ExportReq, _user=Depends(current_user)):
    template = _template_or_404(req.template_id)
    decoded = _decode_images(req.images)
    labels = {_image_key(image): safe_filename(str(image.chart_id or ""))[:80] for image in req.images}
    out = io.BytesIO()
    with zipfile.ZipFile(out, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        for key in sorted(decoded, key=_parse_key):
            page_index, position = _parse_key(key)
            token = labels.get(key) or f"chart_{position}"
            archive.writestr(f"slide_{page_index + 1:02d}_chart_{position}_{token}.png", decoded[key])
    filename = safe_filename(f"{template.get('name') or 'template_report'}_chart_images.zip")
    return Response(
        content=out.getvalue(),
        media_type="application/zip",
        headers={"Content-Disposition": _download_header(filename)},
    )


def _parse_key(key: str) -> tuple[int, int]:
    page_text, _, position_text = str(key or "").partition(":")
    try:
        return int(page_text), int(position_text)
    except ValueError:
        return 0, 0
