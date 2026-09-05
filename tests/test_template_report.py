import base64
import io
import json
import zipfile

import pytest
from PIL import Image
from pptx import Presentation
from fastapi import HTTPException

from core.utils import jsonl_append
from routers import filebrowser, template_report


DEFINITION = """Q1
TABLE = ET
PRODUCT = PRODA
SQL = SELECT tkout_time, root_lot_id, wafer_id, value
RECENT_DAYS = 7
DATE_COLUMN = tkout_time

CHART
TYPE = scatter
X = tkout_time
Y = value
COLOR = root_lot_id
HIGHLIGHT = true
WIDTH = 1200
HEIGHT = 650
MAX_ROWS = 500
"""


@pytest.fixture(autouse=True)
def _isolated_template_report_background(tmp_path, monkeypatch):
    """운영 data_root의 실제 기본 배경이 PPTX 테스트 결과에 섞이지 않게 한다."""
    monkeypatch.setattr(template_report, "SETTINGS_FILE", tmp_path / "template_report_settings.json")
    monkeypatch.setattr(template_report, "BACKGROUND_FILE", tmp_path / "template_report_background.png")


def _seed_chart_history(tmp_path, monkeypatch):
    history_path = tmp_path / "chart_builder_history.jsonl"
    monkeypatch.setattr(filebrowser, "_chart_builder_history_path", lambda: history_path)
    jsonl_append(history_path, {
        "event": "history",
        "history_id": "chart_demo_001",
        "name": "Weekly ET trend",
        "timestamp": "2026-08-10T01:02:03+00:00",
        "username": "engineer",
        "definition_code": DEFINITION,
        "chart": {"type": "scatter", "x": "tkout_time", "y": "value"},
        "sources": [{"id": "q1", "root": "ET", "product": "PRODA"}],
        "row_count": 12,
    })
    return history_path


def _save_demo_template(tmp_path, monkeypatch):
    monkeypatch.setattr(template_report, "STORE_FILE", tmp_path / "template_reports.json")
    _seed_chart_history(tmp_path, monkeypatch)
    result = template_report.save_template(
        template_report.TemplateSaveReq(
            name="Weekly ET report",
            pages=[template_report.TemplatePageReq(
                title="ET trend",
                slots=[template_report.TemplateSlotReq(position=1, chart_id="chart_demo_001")],
            )],
        ),
        {"username": "engineer", "role": "user"},
    )
    return result["template"]


def _png_data_url():
    stream = io.BytesIO()
    Image.new("RGB", (1200, 650), (238, 244, 255)).save(stream, format="PNG")
    return "data:image/png;base64," + base64.b64encode(stream.getvalue()).decode("ascii")


def test_template_run_replays_the_time_window_saved_in_the_chart_code(tmp_path, monkeypatch):
    """시간 창은 차트생성 코드가 가진다 — 보고서는 덧씌우지 않고 그대로 다시 실행한다."""
    saved = _save_demo_template(tmp_path, monkeypatch)

    assert saved["slide_size"] == "LAYOUT_WIDE"
    assert saved["pages"][0]["slots"][0]["definition_code"] == DEFINITION
    # 실행 전에도 "이 보고서가 볼 구간"을 화면에 그대로 보여줄 수 있어야 한다.
    assert saved["pages"][0]["slots"][0]["recent_days"] == 7
    assert saved["pages"][0]["slots"][0]["date_column"] == "tkout_time"

    run = template_report.prepare_run(
        template_report.TemplateRunReq(template_id=saved["id"]),
        {"username": "viewer", "role": "user"},
    )

    request = run["charts"][0]["request"]
    assert run["charts"][0]["chart_width"] == 1200
    assert run["charts"][0]["chart_height"] == 650
    assert request["save_history"] is False
    assert request["sources"][0]["runtime_recent_days"] == 7
    assert request["sources"][0]["runtime_date_column"] == "tkout_time"
    assert request["chart"]["width"] == 1200
    assert request["chart"]["height"] == 650


def test_template_slot_can_override_and_reuse_the_chart_generation_code(tmp_path, monkeypatch):
    monkeypatch.setattr(template_report, "STORE_FILE", tmp_path / "template_reports.json")
    _seed_chart_history(tmp_path, monkeypatch)
    edited = DEFINITION.replace("TYPE = scatter", "TYPE = line").replace("COLOR = root_lot_id", "COLOR = wafer_id")

    result = template_report.save_template(
        template_report.TemplateSaveReq(
            name="Editable chart formula",
            pages=[template_report.TemplatePageReq(
                title="Edited trend",
                slots=[template_report.TemplateSlotReq(
                    position=1,
                    chart_id="chart_demo_001",
                    definition_code=edited,
                )],
            )],
        ),
        {"username": "engineer", "role": "user"},
    )

    saved = result["template"]
    assert saved["pages"][0]["slots"][0]["definition_code"] == edited
    run = template_report.prepare_run(
        template_report.TemplateRunReq(template_id=saved["id"]),
        {"username": "viewer", "role": "user"},
    )
    assert run["charts"][0]["request"]["chart"]["type"] == "line"


def test_available_charts_include_copyable_generation_code(tmp_path, monkeypatch):
    _seed_chart_history(tmp_path, monkeypatch)

    charts = template_report.list_charts({"username": "viewer", "role": "user"})["charts"]

    assert charts[0]["id"] == "chart_demo_001"
    assert charts[0]["definition_code"] == DEFINITION


def test_full_template_code_round_trip_includes_inline_chart_definition(tmp_path, monkeypatch):
    monkeypatch.setattr(template_report, "STORE_FILE", tmp_path / "template_reports.json")
    _seed_chart_history(tmp_path, monkeypatch)
    code = json.dumps({
        "$schema": "flow-template-report/v1",
        "id": "",
        "name": "AI generated process review",
        "options": {"cover": False, "footer": True, "subtitle": "", "repeat_variable": "LOT"},
        "variables": [{"name": "LOT", "label": "Root Lot", "default": "ADV-AB-001"}],
        "pages": [{
            "id": "page_1",
            "title": "Process review",
            "subtitle": "",
            "slots": [{
                "position": 1,
                "kind": "chart",
                "chart_id": "",
                "chart_name": "AI inline trend",
                "x": 4,
                "y": 10,
                "chart_width": 1200,
                "chart_height": 650,
                "definition_code": DEFINITION,
            }],
        }],
    }, ensure_ascii=False)

    parsed = template_report.parse_template_code(
        template_report.TemplateCodeReq(code=code),
        {"username": "engineer", "role": "user"},
    )["template"]

    slot = parsed["pages"][0]["slots"][0]
    assert parsed["name"] == "AI generated process review"
    assert slot["chart_id"] == "inline_p1_1"
    assert slot["chart_name"] == "AI inline trend"
    assert slot["definition_code"] == DEFINITION


def test_full_template_code_reports_json_line_errors():
    with pytest.raises(HTTPException) as exc_info:
        template_report.parse_template_code(
            template_report.TemplateCodeReq(code='{"name": "broken",\n"pages": [}'),
            {"username": "engineer", "role": "user"},
        )

    assert exc_info.value.status_code == 400
    assert "JSON" in str(exc_info.value.detail)


def test_template_ai_assistant_keeps_valid_code_when_company_llm_is_unavailable(tmp_path, monkeypatch):
    from core import llm_adapter

    monkeypatch.setattr(template_report, "STORE_FILE", tmp_path / "template_reports.json")
    _seed_chart_history(tmp_path, monkeypatch)
    monkeypatch.setattr(llm_adapter, "is_available", lambda: False)
    code = json.dumps({
        "name": "AI ready template",
        "pages": [{"title": "Review", "slots": [{
            "position": 1,
            "kind": "chart",
            "chart_id": "chart_demo_001",
            "chart_name": "Weekly ET trend",
            "definition_code": DEFINITION,
        }]}],
    })

    result = template_report.template_assistant(
        template_report.TemplateAssistantReq(
            instruction="trend와 box plot을 한 장에 배치해줘",
            template_code=code,
        ),
        {"username": "engineer", "role": "admin"},
    )

    assert result["changed"] is False
    assert result["llm"]["available"] is False
    assert result["template"]["pages"][0]["slots"][0]["definition_code"] == DEFINITION


def test_blank_page_subtitle_defaults_to_run_date_and_username(tmp_path, monkeypatch):
    saved = _save_demo_template(tmp_path, monkeypatch)
    monkeypatch.setattr(
        template_report,
        "_default_page_subtitle",
        lambda username: f"2026-08-20 {username}",
    )

    run = template_report.prepare_run(
        template_report.TemplateRunReq(template_id=saved["id"]),
        {"username": "viewer", "role": "user"},
    )

    assert run["deck"]["pages"][0]["subtitle"] == "2026-08-20 viewer"


def test_template_run_reuses_tkout_time_highlight_rule_from_chart_code():
    definition = DEFINITION.replace(
        "COLOR = root_lot_id",
        "COLOR = custom\nCOLOR_RULE = tkout_time WITHIN 7 DAYS THEN #ef4444\nCOLOR_ELSE = #cbd5e1",
    )

    request = template_report._chart_request(definition, "테스트 차트")

    assert request["chart"]["color"] == "custom"
    assert request["chart"]["color_rules"] == ["tkout_time WITHIN 7 DAYS THEN #ef4444"]
    assert request["chart"]["color_else"] == "#cbd5e1"
    assert request["save_history"] is False


def test_run_request_has_no_report_level_time_controls(tmp_path, monkeypatch):
    """보고서 화면에 기간 입력이 되살아나면 저장된 차트가 자기 구간을 잃는다."""
    saved = _save_split_template(tmp_path, monkeypatch)

    run = template_report.prepare_run(
        template_report.TemplateRunReq(
            template_id=saved["id"],
            bindings={"PRODUCT": "PRODA", "LOT": "A1234", "SPLIT": "S1"},
        ),
        {"username": "viewer", "role": "user"},
    )

    assert "recent_days" not in run and "date_column" not in run
    # VARIABLE_DEFINITION 에는 RECENT_DAYS 가 없다 → 기간 제한 없이 조회한다.
    assert run["charts"][0]["request"]["sources"][0]["runtime_recent_days"] == 0
    assert not hasattr(template_report.TemplateRunReq(template_id="x"), "recent_days")
    assert not hasattr(template_report.ExportReq(template_id="x"), "recent_days")


def test_template_accepts_chart_name_and_stores_canonical_id(tmp_path, monkeypatch):
    monkeypatch.setattr(template_report, "STORE_FILE", tmp_path / "template_reports.json")
    _seed_chart_history(tmp_path, monkeypatch)

    saved = template_report.save_template(
        template_report.TemplateSaveReq(
            name="Name based report",
            pages=[template_report.TemplatePageReq(
                title="ET trend",
                slots=[template_report.TemplateSlotReq(position=1, chart_id="weekly et trend")],
            )],
        ),
        {"username": "engineer", "role": "user"},
    )["template"]

    slot = saved["pages"][0]["slots"][0]
    assert slot["chart_id"] == "chart_demo_001"
    assert slot["chart_name"] == "Weekly ET trend"


def test_template_names_are_unique_and_name_can_resolve_run(tmp_path, monkeypatch):
    first = _save_demo_template(tmp_path, monkeypatch)
    second = template_report.save_template(
        template_report.TemplateSaveReq(
            name="Weekly ET report",
            pages=[template_report.TemplatePageReq(
                title="Second",
                slots=[template_report.TemplateSlotReq(position=1, chart_id="chart_demo_001")],
            )],
        ),
        {"username": "engineer", "role": "user"},
    )["template"]

    assert first["id"] != second["id"]
    assert first["name"] == "Weekly ET report"
    assert second["name"] == "Weekly ET report (2)"
    run = template_report.prepare_run(
        template_report.TemplateRunReq(template_id="weekly et report (2)"),
        {"username": "viewer", "role": "user"},
    )
    assert run["template"]["id"] == second["id"]


def test_template_pptx_is_wide_and_places_chart_image(tmp_path, monkeypatch):
    saved = _save_demo_template(tmp_path, monkeypatch)
    monkeypatch.setattr(
        template_report,
        "_default_page_subtitle",
        lambda username: f"2026-08-20 {username}",
    )
    request = template_report.ExportReq(
        template_id=saved["id"],
        images=[template_report.ExportImageReq(
            page_index=0,
            position=1,
            chart_id="chart_demo_001",
            data_url=_png_data_url(),
        )],
    )

    response = template_report.export_pptx(request, {"username": "viewer", "role": "user"})
    deck = Presentation(io.BytesIO(response.body))

    assert response.media_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    assert round(deck.slide_width / deck.slide_height, 3) == round(16 / 9, 3)
    # 표지 + 내용 1장 + 맨 뒤 재생성 Appendix
    assert len(deck.slides) == 3
    cover_text = " ".join(getattr(shape, "text", "") for shape in deck.slides[0].shapes)
    assert "Weekly ET report" in cover_text
    content = deck.slides[1]
    assert any(getattr(shape, "text", "") == "ET trend" for shape in content.shapes)
    assert any(getattr(shape, "text", "") == "2026-08-20 viewer" for shape in content.shapes)
    assert len([shape for shape in content.shapes if shape.shape_type == 13]) == 1
    assert "Weekly ET report" not in [getattr(shape, "text", "") for shape in content.shapes]
    appendix_text = " ".join(getattr(shape, "text", "") for shape in deck.slides[-1].shapes)
    assert "Template Report Info" in appendix_text
    assert saved["id"] in appendix_text
    assert "engineer" in appendix_text
    assert "viewer" in appendix_text
    appendix_tables = [shape.table for shape in deck.slides[-1].shapes if shape.has_table]
    assert appendix_tables
    assert appendix_tables[0].cell(1, 2).text == "chart_demo_001"


def test_pasted_background_is_saved_and_used_on_every_pptx_slide(tmp_path, monkeypatch):
    saved = _save_demo_template(tmp_path, monkeypatch)

    stored = template_report.save_template_report_background(
        template_report.BackgroundSaveReq(data_url=_png_data_url()),
        {"username": "report-manager", "role": "admin"},
    )
    background = stored["settings"]["background"]
    assert background["configured"] is True
    assert background["data_url"].startswith("data:image/png;base64,")
    assert background["updated_by"] == "report-manager"
    with Image.open(template_report.BACKGROUND_FILE) as image:
        assert image.format == "PNG"

    response = template_report.export_pptx(
        template_report.ExportReq(
            template_id=saved["id"],
            images=[template_report.ExportImageReq(
                page_index=0,
                position=1,
                chart_id="chart_demo_001",
                data_url=_png_data_url(),
            )],
        ),
        {"username": "viewer", "role": "user"},
    )
    deck = Presentation(io.BytesIO(response.body))
    picture_counts = [len([shape for shape in slide.shapes if shape.shape_type == 13]) for slide in deck.slides]
    assert picture_counts == [1, 2, 1]
    for slide in deck.slides:
        background_picture = next(shape for shape in slide.shapes if shape.shape_type == 13)
        assert background_picture.width >= deck.slide_width or background_picture.height >= deck.slide_height

    removed = template_report.delete_template_report_background(
        {"username": "report-manager", "role": "admin"},
    )
    assert removed["settings"]["background"]["configured"] is False
    assert not template_report.BACKGROUND_FILE.exists()


def test_template_freeform_layout_is_saved_and_used_by_pptx(tmp_path, monkeypatch):
    monkeypatch.setattr(template_report, "STORE_FILE", tmp_path / "template_reports.json")
    _seed_chart_history(tmp_path, monkeypatch)
    saved = template_report.save_template(
        template_report.TemplateSaveReq(
            name="Freeform report",
            pages=[template_report.TemplatePageReq(
                title="Arbitrary placement",
                slots=[template_report.TemplateSlotReq(
                    position=8, chart_id="chart_demo_001",
                    x=10, y=20, width=45, height=37.333,
                )],
            )],
        ),
        {"username": "engineer", "role": "user"},
    )["template"]

    slot = saved["pages"][0]["slots"][0]
    assert {key: slot[key] for key in ("x", "y", "width", "height", "chart_width", "chart_height")} == {
        "x": 10.0, "y": 20.0, "width": 62.5, "height": 60.185,
        "chart_width": 1200, "chart_height": 650,
    }
    payload = template_report._pptx_bytes(
        saved, {"0:8": base64.b64decode(_png_data_url().split(",", 1)[1])}
    )
    deck = Presentation(io.BytesIO(payload))
    picture = next(shape for shape in deck.slides[1].shapes if shape.shape_type == 13)
    assert picture.left / 914400 == pytest.approx(1.3333333, abs=0.02)
    assert picture.top / 914400 == pytest.approx(1.5, abs=0.02)
    assert picture.width / 914400 == pytest.approx(8.333333, abs=0.02)
    assert picture.height / 914400 == pytest.approx(4.513889, abs=0.02)


def test_template_allows_more_than_four_charts_and_legacy_slots_get_coordinates(tmp_path, monkeypatch):
    monkeypatch.setattr(template_report, "STORE_FILE", tmp_path / "template_reports.json")
    _seed_chart_history(tmp_path, monkeypatch)
    saved = template_report.save_template(
        template_report.TemplateSaveReq(
            name="Five charts",
            pages=[template_report.TemplatePageReq(
                title="All charts",
                slots=[template_report.TemplateSlotReq(position=index, chart_id="chart_demo_001") for index in range(1, 6)],
            )],
        ),
        {"username": "engineer", "role": "user"},
    )["template"]

    assert len(saved["pages"][0]["slots"]) == 5
    first = saved["pages"][0]["slots"][0]
    fifth = saved["pages"][0]["slots"][4]
    assert (first["x"], first["y"]) == (3.4, 13.6)
    assert (fifth["x"], fifth["y"]) == (5.0, 16.0)


def test_template_rejects_chart_outside_slide(tmp_path, monkeypatch):
    monkeypatch.setattr(template_report, "STORE_FILE", tmp_path / "template_reports.json")
    _seed_chart_history(tmp_path, monkeypatch)

    with pytest.raises(HTTPException, match="슬라이드 영역"):
        template_report.save_template(
            template_report.TemplateSaveReq(
                name="Invalid layout",
                pages=[template_report.TemplatePageReq(
                    title="Invalid",
                    slots=[template_report.TemplateSlotReq(
                        position=1, chart_id="chart_demo_001", x=80, y=20, width=30, height=30,
                    )],
                )],
            ),
            {"username": "engineer", "role": "user"},
        )


def test_template_image_export_contains_individual_pngs(tmp_path, monkeypatch):
    saved = _save_demo_template(tmp_path, monkeypatch)
    request = template_report.ExportReq(
        template_id=saved["id"],
        images=[template_report.ExportImageReq(
            page_index=0,
            position=1,
            chart_id="chart_demo_001",
            data_url=_png_data_url(),
        )],
    )

    response = template_report.export_images(request, {"username": "viewer", "role": "user"})
    with zipfile.ZipFile(io.BytesIO(response.body)) as archive:
        names = archive.namelist()
        assert names == ["slide_01_chart_1_chart_demo_001.png"]
        assert archive.read(names[0]).startswith(b"\x89PNG")


VARIABLE_DEFINITION = """Q1
TABLE = ET
PRODUCT = {{PRODUCT}}
SQL = SELECT tkout_time, root_lot_id, wafer_id, value WHERE root_lot_id = '{{LOT}}' AND split_cond = '{{SPLIT}}'

CHART
TYPE = box
X = split_cond
Y = value
WIDTH = 1200
HEIGHT = 650
MAX_ROWS = 500
"""


def _seed_variable_chart(tmp_path, monkeypatch):
    history_path = tmp_path / "chart_builder_history.jsonl"
    monkeypatch.setattr(filebrowser, "_chart_builder_history_path", lambda: history_path)
    jsonl_append(history_path, {
        "event": "history",
        "history_id": "chart_ab_001",
        "name": "Split A/B box",
        "timestamp": "2026-08-10T02:00:00+00:00",
        "username": "engineer",
        "definition_code": VARIABLE_DEFINITION,
        "chart": {"type": "box", "x": "split_cond", "y": "value"},
        "sources": [{"id": "q1", "root": "ET", "product": "{{PRODUCT}}"}],
        "row_count": 40,
    })


def _save_split_template(tmp_path, monkeypatch, *, extra_slots=()):
    monkeypatch.setattr(template_report, "STORE_FILE", tmp_path / "template_reports.json")
    _seed_variable_chart(tmp_path, monkeypatch)
    return template_report.save_template(
        template_report.TemplateSaveReq(
            name="Split A/B report",
            variables=[
                template_report.TemplateVariableReq(name="PRODUCT", label="제품", default="PRODA"),
                template_report.TemplateVariableReq(name="LOT", label="랏"),
                template_report.TemplateVariableReq(name="SPLIT", label="조건"),
            ],
            options=template_report.TemplateOptionsReq(
                subtitle="{{PRODUCT}} · {{LOT}}",
                repeat_variable="LOT",
            ),
            pages=[template_report.TemplatePageReq(
                title="{{LOT}} Split 비교",
                subtitle="{{PRODUCT}}",
                slots=[
                    template_report.TemplateSlotReq(position=1, chart_id="chart_ab_001"),
                    *extra_slots,
                ],
            )],
        ),
        {"username": "engineer", "role": "user"},
    )["template"]


def test_template_variables_are_detected_and_repeat_expands_pages(tmp_path, monkeypatch):
    saved = _save_split_template(tmp_path, monkeypatch)

    assert [item["name"] for item in saved["variables"]] == ["PRODUCT", "LOT", "SPLIT"]
    run = template_report.prepare_run(
        template_report.TemplateRunReq(
            template_id=saved["id"],
            bindings={"PRODUCT": "PRODA", "SPLIT": "S1"},
            repeat_values=["A1234, A5678"],
        ),
        {"username": "viewer", "role": "user"},
    )

    pages = run["deck"]["pages"]
    assert len(pages) == 2
    assert [page["title"] for page in pages] == ["A1234 Split 비교", "A5678 Split 비교"]
    assert [page["subtitle"] for page in pages] == ["PRODA", "PRODA"]
    sql_first = run["charts"][0]["request"]["sources"][0]["sql"]
    sql_second = run["charts"][1]["request"]["sources"][0]["sql"]
    assert "root_lot_id = 'A1234'" in sql_first and "split_cond = 'S1'" in sql_first
    assert "root_lot_id = 'A5678'" in sql_second
    assert run["charts"][0]["request"]["sources"][0]["product"] == "PRODA"


def test_report_context_updates_all_chart_filters_time_and_color(tmp_path, monkeypatch):
    saved = _save_split_template(tmp_path, monkeypatch)
    context = template_report.TemplateRunContextReq(
        root_lot_ids=["A1234", "A5678"],
        wafer_ids=["1", "2"],
        override_recent_days=True,
        recent_days=14,
        date_column="tkout_time",
        color_rules=["root_lot_id = 'A1234' AND wafer_id = '1' THEN #dc2626"],
        color_else="#cbd5e1",
    )

    run = template_report.prepare_run(
        template_report.TemplateRunReq(
            template_id=saved["id"],
            bindings={"PRODUCT": "PRODA", "LOT": "A1234", "SPLIT": "S1"},
            context=context,
        ),
        {"username": "viewer", "role": "user"},
    )

    request = run["charts"][0]["request"]
    source = request["sources"][0]
    assert source["runtime_root_lot_ids"] == ["A1234", "A5678"]
    assert source["runtime_wafer_ids"] == ["1", "2"]
    assert source["runtime_lot_wafer_pairs"] == [{"root_lot_id": "A1234", "wafer_id": "1"}]
    assert source["runtime_recent_days"] == 14
    assert source["runtime_date_column"] == "tkout_time"
    assert request["chart"]["color"] == "custom"
    assert request["chart"]["color_rules"] == context.color_rules


def test_slot_renders_once_at_its_declared_size(tmp_path, monkeypatch):
    """조건 비교는 ChartBuilder 코드가 만든다 — 템플릿은 슬롯을 복제하지 않는다."""
    saved = _save_split_template(tmp_path, monkeypatch)

    run = template_report.prepare_run(
        template_report.TemplateRunReq(
            template_id=saved["id"],
            bindings={"PRODUCT": "PRODA", "SPLIT": "S1"},
            repeat_values=["A1234"],
        ),
        {"username": "viewer", "role": "user"},
    )

    blocks = run["deck"]["pages"][0]["blocks"]
    assert [block["key"] for block in blocks] == ["0:1"]
    assert blocks[0]["chart_width"] == 1200
    assert "split_cond = 'S1'" in blocks[0]["request"]["sources"][0]["sql"]


def test_missing_variable_value_stops_the_run(tmp_path, monkeypatch):
    saved = _save_split_template(tmp_path, monkeypatch)

    with pytest.raises(HTTPException) as excinfo:
        template_report.prepare_run(
            template_report.TemplateRunReq(template_id=saved["id"], bindings={"PRODUCT": "PRODA"}),
            {"username": "viewer", "role": "user"},
        )
    assert "LOT" in str(excinfo.value.detail)


def test_variable_value_cannot_break_out_of_the_sql_literal(tmp_path, monkeypatch):
    saved = _save_split_template(tmp_path, monkeypatch)

    with pytest.raises(HTTPException) as excinfo:
        template_report.prepare_run(
            template_report.TemplateRunReq(
                template_id=saved["id"],
                bindings={"PRODUCT": "PRODA", "SPLIT": "S1", "LOT": "A1' OR 1=1 --"},
            ),
            {"username": "viewer", "role": "user"},
        )
    assert "쓸 수 없는 문자" in str(excinfo.value.detail)


def test_pptx_has_compact_header_cover_and_table_blocks(tmp_path, monkeypatch):
    saved = _save_split_template(tmp_path, monkeypatch, extra_slots=(
        template_report.TemplateSlotReq(
            position=2, kind="split", product="{{PRODUCT}}", lot="{{LOT}}",
            x=5, y=62, width=60, height=30, title="Split 조건",
        ),
        template_report.TemplateSlotReq(
            position=3, kind="text", text="{{LOT}} 판정: 이상 없음",
            x=68, y=62, width=28, height=20,
        ),
    ))

    request = template_report.ExportReq(
        template_id=saved["id"],
        bindings={"PRODUCT": "PRODA", "LOT": "A1234", "SPLIT": "S1"},
        images=[template_report.ExportImageReq(key="0:1", chart_id="chart_ab_001", data_url=_png_data_url())],
        tables=[template_report.ExportTableReq(
            key="0:2",
            columns=["parameter", "#1", "#2"],
            rows=[["KNOB_Vt_Split", "S1", "S2"], ["KNOB_Dose", "1.0", "1.2"]],
            note="2 params",
        )],
    )

    response = template_report.export_pptx(request, {"username": "viewer", "role": "user"})
    deck = Presentation(io.BytesIO(response.body))

    assert len(deck.slides) == 3
    content = deck.slides[1]
    texts = [getattr(shape, "text", "") for shape in content.shapes]
    assert "A1234 Split 비교" in texts
    compact_title = next(shape for shape in content.shapes if getattr(shape, "text", "") == "A1234 Split 비교")
    assert compact_title.top < template_report.TITLE_BAR_HEIGHT_IN * 914400
    assert compact_title.text_frame.paragraphs[0].font.size.pt == 14
    assert "A1234 판정: 이상 없음" in texts
    assert "1 / 1" in texts
    tables = [shape for shape in content.shapes if shape.has_table]
    assert len(tables) == 1
    assert tables[0].table.cell(1, 0).text == "KNOB_Vt_Split"


def test_stats_block_points_at_its_source_chart_and_accepts_legacy_kind(tmp_path, monkeypatch):
    saved = _save_split_template(tmp_path, monkeypatch, extra_slots=(
        # 예전 템플릿이 쓰던 kind — 지금은 통계표 하나로 합쳐 승격시킨다.
        template_report.TemplateSlotReq(
            position=4, kind="ab_stats", source_position=1,
            x=5, y=64, width=52, height=28,
        ),
    ))

    assert saved["pages"][0]["slots"][1]["kind"] == "stats"
    run = template_report.prepare_run(
        template_report.TemplateRunReq(
            template_id=saved["id"],
            bindings={"PRODUCT": "PRODA", "LOT": "A1234", "SPLIT": "S1"},
        ),
        {"username": "viewer", "role": "user"},
    )

    stats = next(block for block in run["deck"]["pages"][0]["blocks"] if block["kind"] == "stats")
    assert stats["source_key"] == "0:1"
    assert stats["stats"] == "n,mean,median,std"


def test_common_legend_block_uses_chart_groups_and_exports_as_native_shapes(tmp_path, monkeypatch):
    monkeypatch.setattr(template_report, "STORE_FILE", tmp_path / "template_reports.json")
    _seed_chart_history(tmp_path, monkeypatch)
    saved = template_report.save_template(
        template_report.TemplateSaveReq(
            name="Shared legend report",
            pages=[template_report.TemplatePageReq(
                title="Common legend",
                slots=[
                    template_report.TemplateSlotReq(position=1, chart_id="chart_demo_001"),
                    template_report.TemplateSlotReq(
                        position=2, kind="legend", source_position=1,
                        x=70, y=8, width=26, height=11, title="Root Lot",
                    ),
                ],
            )],
        ),
        {"username": "engineer", "role": "user"},
    )["template"]

    run = template_report.prepare_run(
        template_report.TemplateRunReq(template_id=saved["id"]),
        {"username": "viewer", "role": "user"},
    )
    legend = next(block for block in run["deck"]["pages"][0]["blocks"] if block["kind"] == "legend")
    assert legend["source_key"] == "0:1"

    response = template_report.export_pptx(
        template_report.ExportReq(
            template_id=saved["id"],
            images=[template_report.ExportImageReq(key="0:1", chart_id="chart_demo_001", data_url=_png_data_url())],
            tables=[template_report.ExportTableReq(
                key="0:2", title="Root Lot", columns=["Label", "Color", "Count"],
                rows=[["LOT-A", "#6366f1", "12"], ["LOT-B", "#f59e0b", "9"]],
            )],
        ),
        {"username": "viewer", "role": "user"},
    )
    deck = Presentation(io.BytesIO(response.body))
    content_text = " ".join(getattr(shape, "text", "") for shape in deck.slides[1].shapes)
    assert "Root Lot" in content_text
    assert "LOT-A (12)" in content_text
    assert not any(shape.has_table for shape in deck.slides[1].shapes)


def test_download_header_supports_korean_template_names():
    header = template_report._download_header("주간 ET 분석_7days.pptx")

    header.encode("latin-1")
    assert "filename*=UTF-8''" in header
    assert header.endswith(".pptx")


def test_appendix_time_keeps_local_clock_without_timezone_label():
    shown = template_report._format_report_time("2026-08-21T02:15:00+00:00")

    assert len(shown) == 16
    assert shown[4] == "-" and shown[7] == "-" and shown[10] == " " and shown[13] == ":"
    assert "표준시" not in shown and "KST" not in shown and "UTC" not in shown


def test_compact_appendix_fits_eighteen_chart_rows_on_first_info_page():
    deck = Presentation()
    deck.slide_width = int(template_report.SLIDE_WIDTH_IN * 914400)
    deck.slide_height = int(template_report.SLIDE_HEIGHT_IN * 914400)
    template = {
        "id": "report_tpl_dense",
        "name": "Dense chart inventory",
        "pages": [{
            "title": "Dense",
            "slots": [
                {
                    "position": position, "kind": "chart", "chart_id": f"chart_{position:02d}",
                    "chart_name": f"Chart {position:02d}", "definition_code": DEFINITION,
                    "x": 4, "y": 12, "chart_width": 870, "chart_height": 390,
                }
                for position in range(1, 19)
            ],
        }],
        "options": {"cover": False, "footer": True},
        "created_by": "author", "created_at": "2026-08-20T00:00:00+00:00",
        "updated_by": "author", "updated_at": "2026-08-21T00:00:00+00:00",
    }
    template_report._add_template_info_appendix(
        deck, template, {"bindings": {}, "repeat_variable": "", "repeat_values": [], "context": {}},
        report_user="viewer", report_generated_at="2026-08-21T02:15:00+00:00",
    )

    assert len(deck.slides) == 1
    appendix = deck.slides[0]
    inventory = next(shape for shape in appendix.shapes if shape.has_table)
    footer = next(shape for shape in appendix.shapes if getattr(shape, "text", "") == "Appendix 1 / 1")
    assert len(inventory.table.rows) == 19  # header + 18 charts
    assert inventory.top + inventory.height < footer.top
